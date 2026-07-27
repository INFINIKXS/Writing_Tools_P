"""
Document conversion API routes and helper tools.
"""
import io
import os
import shutil
import asyncio
import subprocess
import tempfile
import uuid
import zipfile
import json
from enum import Enum
from concurrent.futures import ThreadPoolExecutor
from pathlib import Path
from typing import List, Optional, Tuple, Union

from fastapi import APIRouter, BackgroundTasks, UploadFile, File, Form, HTTPException
from fastapi.responses import JSONResponse, FileResponse
from starlette.responses import StreamingResponse
from PyPDF2 import PdfReader, PdfWriter
from docx import Document

from core.config import (
    PDF2DOCX_AVAILABLE, Pdf2DocxConverter,
    PILLOW_AVAILABLE, Image,
    PDF2IMAGE_AVAILABLE, convert_from_bytes,
    TESSERACT_AVAILABLE, pytesseract,
    POPPLER_PATH,
    PIKEPDF_AVAILABLE, pikepdf,
    PDF2PPTX_AVAILABLE, pdf2pptx,
    PDF_TO_EXCEL_AVAILABLE,
    WEASYPRINT_AVAILABLE, weasyprint, _find_ghostscript,
)
from utils.text_extraction import extract_pdf_text

router = APIRouter()


# ─── Universal Async Job Store ──────────────────────────────────────────────
class JobStatus(str, Enum):
    PENDING = "pending"
    PROCESSING = "processing"
    DONE = "done"
    ERROR = "error"

_job_store: dict = {}  # job_id -> {status, result_path, tmp_dir, filename, content_type, error}
_executor = ThreadPoolExecutor(max_workers=4)
JOB_TTL_SECONDS = 1800  # 30 minutes — jobs survive for re-download


async def _schedule_job_cleanup(job_id: str):
    """Delete job entry and temp files after TTL expires."""
    await asyncio.sleep(JOB_TTL_SECONDS)
    job = _job_store.pop(job_id, None)
    if job and job.get("tmp_dir"):
        shutil.rmtree(job["tmp_dir"], ignore_errors=True)


def _run_word_to_pdf_sync(job_id: str, in_path: str, tmp_dir: str, soffice: str):
    """Blocking LibreOffice conversion — runs in thread pool, no timeout ceiling."""
    try:
        _job_store[job_id]["status"] = JobStatus.PROCESSING
        result = subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf", "--outdir", tmp_dir, in_path],
            capture_output=True,
            text=True,
        )
        if result.returncode != 0:
            _job_store[job_id]["status"] = JobStatus.ERROR
            _job_store[job_id]["error"] = f"LibreOffice conversion failed: {result.stderr}"
            return

        # LibreOffice names the output after the input file stem
        pdf_path = os.path.join(tmp_dir, Path(in_path).stem + ".pdf")
        if not os.path.isfile(pdf_path):
            pdf_files = [f for f in os.listdir(tmp_dir) if f.endswith(".pdf")]
            if pdf_files:
                pdf_path = os.path.join(tmp_dir, pdf_files[0])
            else:
                _job_store[job_id]["status"] = JobStatus.ERROR
                _job_store[job_id]["error"] = "LibreOffice did not produce a PDF output."
                return

        _job_store[job_id]["status"] = JobStatus.DONE
        _job_store[job_id]["result_path"] = pdf_path
    except Exception as e:
        _job_store[job_id]["status"] = JobStatus.ERROR
        _job_store[job_id]["error"] = str(e)


def _run_pdf_to_word_sync(job_id: str, in_path: str, tmp_dir: str):
    """Blocking PDF→DOCX conversion — runs in thread pool."""
    try:
        _job_store[job_id]["status"] = JobStatus.PROCESSING
        docx_path = os.path.join(tmp_dir, "output.docx")

        cv = Pdf2DocxConverter(in_path)
        cv.convert(docx_path)
        cv.close()

        # Fall back to OCR if very little text was extracted
        doc = Document(docx_path)
        text_content = "\n".join(p.text for p in doc.paragraphs).strip()
        if len(text_content) < 50 and TESSERACT_AVAILABLE and PDF2IMAGE_AVAILABLE:
            with open(in_path, "rb") as f:
                pdf_bytes = f.read()
            _ocr_pdf_to_docx(pdf_bytes, docx_path)

        _job_store[job_id]["status"] = JobStatus.DONE
        _job_store[job_id]["result_path"] = docx_path
    except Exception as e:
        _job_store[job_id]["status"] = JobStatus.ERROR
        _job_store[job_id]["error"] = str(e)


def _run_pdf_to_text_sync(job_id: str, in_path: str, tmp_dir: str, stem: str):
    """Blocking PDF→TXT extraction — runs in thread pool."""
    try:
        _job_store[job_id]["status"] = JobStatus.PROCESSING
        with open(in_path, "rb") as f:
            file_bytes = f.read()

        text = extract_pdf_text(file_bytes)
        if len(text.strip()) < 50 and TESSERACT_AVAILABLE and PDF2IMAGE_AVAILABLE:
            ocr_text = _ocr_pdf_to_text(file_bytes)
            if ocr_text.strip():
                text = ocr_text

        out_path = os.path.join(tmp_dir, f"{stem}.txt")
        with open(out_path, "w", encoding="utf-8") as f:
            f.write(text)

        _job_store[job_id]["status"] = JobStatus.DONE
        _job_store[job_id]["result_path"] = out_path
    except Exception as e:
        _job_store[job_id]["status"] = JobStatus.ERROR
        _job_store[job_id]["error"] = str(e)


def _run_pdf_to_images_sync(job_id: str, in_paths: list, tmp_dir: str, stems: list):
    """Blocking PDF→images conversion for one or more PDFs — runs in thread pool.
    
    All images from all PDFs are packed into a single ZIP with a 100% flat structure.
    No nested sub-folders are created in the ZIP archive.
    """
    try:
        _job_store[job_id]["status"] = JobStatus.PROCESSING

        used_names = set()
        all_entries = []  # list of (arcname, jpeg_bytes)

        for in_path, stem in zip(in_paths, stems):
            clean_stem = Path(stem).name
            with open(in_path, "rb") as f:
                file_bytes = f.read()
            images = convert_from_bytes(file_bytes, dpi=200, fmt="jpeg", poppler_path=POPPLER_PATH)
            if not images:
                raise ValueError(f"No pages found in '{clean_stem}.pdf'.")
            for i, img in enumerate(images, 1):
                img_buf = io.BytesIO()
                img.save(img_buf, "JPEG", quality=90)
                img_buf.seek(0)

                base_arcname = f"{clean_stem}_page_{i}.jpg"
                arcname = base_arcname
                counter = 1
                while arcname in used_names:
                    arcname = f"{clean_stem}_page_{i}_{counter}.jpg"
                    counter += 1
                used_names.add(arcname)

                all_entries.append((arcname, img_buf.read()))

        if len(in_paths) == 1 and len(all_entries) == 1:
            # Single PDF, single page → return bare JPG
            arcname, jpeg_bytes = all_entries[0]
            out_path = os.path.join(tmp_dir, arcname)
            with open(out_path, "wb") as fout:
                fout.write(jpeg_bytes)
            _job_store[job_id]["result_path"] = out_path
            _job_store[job_id]["filename"] = arcname
            _job_store[job_id]["content_type"] = "image/jpeg"
        else:
            zip_name = f"{Path(stems[0]).name}_images.zip" if len(in_paths) == 1 else "pdf_images.zip"
            zip_path = os.path.join(tmp_dir, zip_name)
            with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
                for arcname, jpeg_bytes in all_entries:
                    zf.writestr(arcname, jpeg_bytes)
            _job_store[job_id]["result_path"] = zip_path
            _job_store[job_id]["filename"] = zip_name
            _job_store[job_id]["content_type"] = "application/zip"

        _job_store[job_id]["status"] = JobStatus.DONE
    except Exception as e:
        error_msg = str(e)
        if "poppler" in error_msg.lower() or "pdftoppm" in error_msg.lower():
            error_msg = "Poppler is not installed or not on PATH."
        _job_store[job_id]["status"] = JobStatus.ERROR
        _job_store[job_id]["error"] = error_msg


def _parse_page_ranges(ranges_str: str, total_pages: int) -> list:
    """Parse a range string like '1-3, 5, 7-9' into sorted 0-indexed page indices."""
    if not ranges_str or ranges_str.strip().lower() in ("all", "all pages", "1-all", "1-end", "1-"):
        return list(range(total_pages))
    pages = set()
    for part in ranges_str.split(','):
        part = part.strip()
        if not part:
            continue
        if '-' in part:
            try:
                start_s, end_s = part.split('-', 1)
                start_s = start_s.strip()
                end_s = end_s.strip()
                start = 0 if not start_s or start_s.lower() == 'start' else int(start_s) - 1
                if end_s.lower() in ('end', 'max', 'all', ''):
                    end = total_pages - 1
                else:
                    end = int(end_s) - 1
                for p in range(max(0, start), min(end + 1, total_pages)):
                    pages.add(p)
            except ValueError:
                pass
        else:
            try:
                if part.lower() in ('all', 'end', 'max'):
                    for p in range(total_pages):
                        pages.add(p)
                else:
                    p = int(part.strip()) - 1
                    if 0 <= p < total_pages:
                        pages.add(p)
            except ValueError:
                pass
    return sorted(pages)


def _parse_multi_pdf_order(order_str: str, doc_page_counts: List[int]) -> List[Tuple[int, int]]:
    """
    Parse an order string for organizing one or multiple PDFs.
    Supports formats like:
    - '0:0, 0:1, 1:0' (0-indexed or 1-indexed file_index:page_index)
    - '3, 1, 2' or '1-3, 5' (1-indexed page sequence for file 0)
    Returns a list of (file_index, page_index) tuples with 0-based indices, preserving exact sequence order.
    """
    if not doc_page_counts:
        return []

    if not order_str or not order_str.strip():
        return [(f_idx, p_idx) for f_idx, count in enumerate(doc_page_counts) for p_idx in range(count)]

    tokens = [t.strip() for t in order_str.split(',') if t.strip()]
    if not tokens:
        return [(f_idx, p_idx) for f_idx, count in enumerate(doc_page_counts) for p_idx in range(count)]

    has_colon = any(':' in t for t in tokens)

    # First pass: check if 0-based page indexing is used when colons are present (e.g. presence of ':0')
    is_zero_based = False
    for t in tokens:
        if ':' in t:
            _, p_spec = t.split(':', 1)
            p_spec = p_spec.strip()
            if p_spec == '0' or p_spec.startswith('0-') or p_spec.endswith('-0') or '-0-' in p_spec:
                is_zero_based = True
                break

    result = []
    for t in tokens:
        if ':' in t:
            f_str, p_str = t.split(':', 1)
            try:
                f_idx = int(f_str.strip())
            except ValueError:
                continue
            p_str = p_str.strip()
        else:
            f_idx = 0
            p_str = t

        if f_idx < 0 or f_idx >= len(doc_page_counts):
            continue

        total_p = doc_page_counts[f_idx]

        if '-' in p_str and not p_str.startswith('-'):
            parts = p_str.split('-', 1)
            try:
                start = int(parts[0].strip())
                end = int(parts[1].strip())
                if has_colon:
                    if is_zero_based:
                        step = 1 if start <= end else -1
                        p_indices = list(range(start, end + step, step))
                    else:
                        start_0 = start - 1
                        end_0 = end - 1
                        step = 1 if start_0 <= end_0 else -1
                        p_indices = list(range(start_0, end_0 + step, step))
                else:
                    start_0 = start - 1
                    end_0 = end - 1
                    step = 1 if start_0 <= end_0 else -1
                    p_indices = list(range(start_0, end_0 + step, step))
            except ValueError:
                continue
        else:
            try:
                p_val = int(p_str)
                if has_colon:
                    if is_zero_based:
                        p_indices = [p_val]
                    else:
                        if 1 <= p_val <= total_p:
                            p_indices = [p_val - 1]
                        else:
                            p_indices = [p_val]
                else:
                    if 1 <= p_val <= total_p:
                        p_indices = [p_val - 1]
                    else:
                        p_indices = [p_val]
            except ValueError:
                continue

        for p in p_indices:
            if 0 <= p < total_p:
                result.append((f_idx, p))

    if not result:
        return [(f_idx, p_idx) for f_idx, count in enumerate(doc_page_counts) for p_idx in range(count)]

    return result


def _run_merge_pdfs_sync(job_id: str, file_paths: list, tmp_dir: str):
    """Merge multiple PDFs — runs in thread pool."""
    try:
        _job_store[job_id]["status"] = JobStatus.PROCESSING
        out_path = os.path.join(tmp_dir, "merged_output.pdf")
        try:
            import fitz
            doc_out = fitz.open()
            for path in file_paths:
                with fitz.open(path) as doc_in:
                    doc_out.insert_pdf(doc_in)
            doc_out.save(out_path)
            doc_out.close()
        except Exception:
            writer = PdfWriter()
            open_handles = []
            try:
                for path in file_paths:
                    f = open(path, 'rb')
                    open_handles.append(f)
                    reader = PdfReader(f)
                    for page in reader.pages:
                        writer.add_page(page)
                with open(out_path, 'wb') as f_out:
                    writer.write(f_out)
            finally:
                for h in open_handles:
                    try:
                        h.close()
                    except Exception:
                        pass
        _job_store[job_id]["status"] = JobStatus.DONE
        _job_store[job_id]["result_path"] = out_path
    except Exception as e:
        _job_store[job_id]["status"] = JobStatus.ERROR
        _job_store[job_id]["error"] = str(e)



def _run_compress_pdf_sync(job_id: str, in_path: str, tmp_dir: str, compression_level: str = "recommended", custom_dpi: int = 0, custom_quality: int = 0):
    """Compress a PDF — using PyMuPDF (fitz) with image re-encoding, DPI downsampling, stream deflation, and raster fallback."""
    try:
        _job_store[job_id]["status"] = JobStatus.PROCESSING
        original_size = os.path.getsize(in_path)
        out_path = os.path.join(tmp_dir, "compressed_output.pdf")

        try:
            import fitz
            from PIL import Image

            # Determine target dimensions & quality based on level/presets
            if compression_level == "extreme":
                max_dim = 950      # ~72-96 DPI
                quality = 42
            elif compression_level == "less":
                max_dim = 2400     # ~200-300 DPI
                quality = 85
            else:  # "recommended"
                max_dim = 1400     # ~120-150 DPI
                quality = 62

            if custom_dpi and custom_dpi > 0:
                max_dim = int((custom_dpi / 72.0) * 612)
            if custom_quality and custom_quality > 0:
                quality = custom_quality

            doc = fitz.open(in_path)

            # Pass 1: Extract, downsample, and re-encode embedded images
            processed_xrefs = set()
            for page in doc:
                for img_info in page.get_images(full=True):
                    xref = img_info[0]
                    if xref in processed_xrefs:
                        continue
                    processed_xrefs.add(xref)
                    try:
                        base_image = doc.extract_image(xref)
                        if not base_image:
                            continue
                        img_bytes = base_image["image"]

                        img = Image.open(io.BytesIO(img_bytes))
                        w, h = img.size
                        max_d = max(w, h)

                        # Downsample if larger than max_dim
                        if max_d > max_dim:
                            scale = max_dim / float(max_d)
                            img = img.resize((max(1, int(w * scale)), max(1, int(h * scale))), Image.Resampling.LANCZOS)

                        # Convert RGBA/P palette to RGB with white background
                        if img.mode in ("RGBA", "P", "LA"):
                            bg = Image.new("RGB", img.size, (255, 255, 255))
                            if img.mode == "P":
                                img = img.convert("RGBA")
                            if "A" in img.mode:
                                bg.paste(img, mask=img.split()[-1])
                            else:
                                bg.paste(img)
                            img = bg
                        elif img.mode != "RGB":
                            img = img.convert("RGB")

                        out_b = io.BytesIO()
                        img.save(out_b, format="JPEG", quality=quality, optimize=True)
                        new_bytes = out_b.getvalue()

                        if len(new_bytes) < len(img_bytes):
                            page.replace_image(xref, stream=new_bytes)
                    except Exception:
                        pass

            doc.save(
                out_path,
                garbage=4,
                deflate=True,
                deflate_images=True,
                deflate_fonts=True,
                clean=True,
            )
            doc.close()

            compressed_size = os.path.getsize(out_path)

            # Pass 2: Rasterization Fallback if Pass 1 yield < 12% reduction on extreme/recommended modes
            # (Essential for Word-converted PDFs with high-res vector objects/paths)
            if (original_size - compressed_size) / float(original_size) < 0.12 and compression_level in ("extreme", "recommended"):
                target_dpi = 100 if compression_level == "extreme" else 130
                target_q = 40 if compression_level == "extreme" else 60

                doc_orig = fitz.open(in_path)
                doc_raster = fitz.open()

                for p in doc_orig:
                    pix = p.get_pixmap(dpi=target_dpi)
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    img_buf = io.BytesIO()
                    img.save(img_buf, format="JPEG", quality=target_q, optimize=True)

                    new_p = doc_raster.new_page(width=p.rect.width, height=p.rect.height)
                    new_p.insert_image(new_p.rect, stream=img_buf.getvalue())

                out_path_pass2 = os.path.join(tmp_dir, "compressed_raster.pdf")
                doc_raster.save(out_path_pass2, garbage=4, deflate=True, deflate_images=True)
                doc_raster.close()
                doc_orig.close()

                size_pass2 = os.path.getsize(out_path_pass2)
                if size_pass2 < compressed_size:
                    out_path = out_path_pass2
                    compressed_size = size_pass2

        except Exception:
            # Fallback to PyPDF2 if fitz fails
            with open(in_path, 'rb') as f:
                file_bytes = f.read()
            reader = PdfReader(io.BytesIO(file_bytes))
            writer = PdfWriter()
            for page in reader.pages:
                page.compress_content_streams()
                writer.add_page(page)
            writer.add_metadata({'/Producer': 'WritingTools', '/Creator': 'WritingTools PDF Compressor'})
            with open(out_path, 'wb') as f:
                writer.write(f)

        compressed_size = os.path.getsize(out_path)
        # If output file turned out larger than original, revert to original
        if compressed_size > original_size:
            shutil.copyfile(in_path, out_path)
            compressed_size = original_size

        ratio = f"{(1 - compressed_size / original_size) * 100:.1f}%" if original_size > 0 else "0%"
        _job_store[job_id]["status"] = JobStatus.DONE
        _job_store[job_id]["result_path"] = out_path
        _job_store[job_id]["extra_headers"] = {
            "X-Original-Size": str(original_size),
            "X-Compressed-Size": str(compressed_size),
            "X-Compression-Ratio": ratio,
        }
    except Exception as e:
        _job_store[job_id]["status"] = JobStatus.ERROR
        _job_store[job_id]["error"] = str(e)


def _run_compress_image_sync(
    job_id: str,
    file_paths: list,
    tmp_dir: str,
    quality: int = 75,
    target_format: str = "original",
    max_dim: int = 0,
):
    """Compress image files synchronously — runs in thread pool."""
    try:
        from PIL import Image, ImageOps
        _job_store[job_id]["status"] = JobStatus.PROCESSING

        original_size = sum(os.path.getsize(p) for p in file_paths if os.path.exists(p))
        target_fmt_lower = (target_format or "original").lower().strip()
        compressed_items = []

        for p in file_paths:
            orig_filename = Path(p).name
            if "___" in orig_filename:
                orig_filename = orig_filename.split("___", 1)[1]
            orig_stem = Path(orig_filename).stem
            orig_ext = Path(orig_filename).suffix.lower()

            if target_fmt_lower in ("jpg", "jpeg"):
                out_fmt = "JPEG"
                out_ext = ".jpg"
            elif target_fmt_lower == "webp":
                out_fmt = "WEBP"
                out_ext = ".webp"
            elif target_fmt_lower == "png":
                out_fmt = "PNG"
                out_ext = ".png"
            else:  # "original"
                if orig_ext in (".jpg", ".jpeg"):
                    out_fmt = "JPEG"
                    out_ext = orig_ext
                elif orig_ext == ".webp":
                    out_fmt = "WEBP"
                    out_ext = ".webp"
                elif orig_ext == ".png":
                    out_fmt = "PNG"
                    out_ext = ".png"
                elif orig_ext == ".bmp":
                    out_fmt = "BMP"
                    out_ext = ".bmp"
                elif orig_ext in (".tiff", ".tif"):
                    out_fmt = "TIFF"
                    out_ext = orig_ext
                else:
                    out_fmt = "JPEG"
                    out_ext = ".jpg"

            img = Image.open(p)
            try:
                img = ImageOps.exif_transpose(img)
            except Exception:
                pass

            w, h = img.size
            if max_dim > 0 and max(w, h) > max_dim:
                scale = float(max_dim) / float(max(w, h))
                new_w = max(1, int(round(w * scale)))
                new_h = max(1, int(round(h * scale)))
                img = img.resize((new_w, new_h), Image.Resampling.LANCZOS)

            dest_file_name = f"proc_{uuid.uuid4().hex[:6]}{out_ext}"
            dest_path = os.path.join(tmp_dir, dest_file_name)

            if out_fmt == "JPEG":
                if img.mode in ("RGBA", "LA", "P", "CMYK", "1"):
                    img = img.convert("RGB")
                img.save(dest_path, format="JPEG", quality=quality, optimize=True, subsampling=2)
            elif out_fmt == "WEBP":
                if img.mode in ("CMYK", "P", "1"):
                    img = img.convert("RGBA" if "A" in img.mode else "RGB")
                img.save(dest_path, format="WEBP", quality=quality, optimize=True, method=6)
            elif out_fmt == "PNG":
                if img.mode not in ("RGB", "RGBA"):
                    if img.mode in ("P", "PA", "LA", "RGBA"):
                        img = img.convert("RGBA")
                    else:
                        img = img.convert("RGB")
                try:
                    colors = img.getcolors(maxcolors=256)
                    if colors is not None and len(colors) <= 256:
                        img = img.quantize(colors=256)
                except Exception:
                    pass
                img.save(dest_path, format="PNG", optimize=True, compress_level=9)
            else:
                img.save(dest_path, format=out_fmt)

            compressed_items.append((dest_path, f"{orig_stem}{out_ext}"))

        if len(file_paths) == 1:
            temp_dest_path, out_filename = compressed_items[0]
            out_ext = Path(out_filename).suffix.lower()
            final_single_path = os.path.join(tmp_dir, f"compressed_output{out_ext}")
            shutil.move(temp_dest_path, final_single_path)

            content_type_map = {
                ".jpg": "image/jpeg",
                ".jpeg": "image/jpeg",
                ".png": "image/png",
                ".webp": "image/webp",
                ".bmp": "image/bmp",
                ".tiff": "image/tiff",
                ".tif": "image/tiff",
            }
            content_type = content_type_map.get(out_ext, "image/jpeg")

            result_path = final_single_path
            _job_store[job_id]["result_path"] = result_path
            _job_store[job_id]["filename"] = out_filename
            _job_store[job_id]["content_type"] = content_type
        else:
            zip_path = os.path.join(tmp_dir, "compressed_images.zip")
            used_names = set()
            with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
                for item_path, ideal_name in compressed_items:
                    stem = Path(ideal_name).stem
                    ext = Path(ideal_name).suffix
                    final_name = ideal_name
                    counter = 1
                    while final_name in used_names:
                        final_name = f"{stem}_{counter}{ext}"
                        counter += 1
                    used_names.add(final_name)
                    zf.write(item_path, arcname=final_name)

            result_path = zip_path
            _job_store[job_id]["result_path"] = result_path
            _job_store[job_id]["filename"] = "compressed_images.zip"
            _job_store[job_id]["content_type"] = "application/zip"

        compressed_size = os.path.getsize(result_path)
        ratio = f"{(1 - compressed_size / original_size) * 100:.1f}%" if original_size > 0 else "0.0%"

        _job_store[job_id]["extra_headers"] = {
            "X-Original-Size": str(original_size),
            "X-Compressed-Size": str(compressed_size),
            "X-Compression-Ratio": ratio,
        }
        _job_store[job_id]["status"] = JobStatus.DONE
    except Exception as e:
        _job_store[job_id]["status"] = JobStatus.ERROR
        _job_store[job_id]["error"] = str(e)


def _run_image_to_pdf_sync(job_id: str, file_paths: list, tmp_dir: str):
    """Convert images to a single PDF — runs in thread pool."""
    try:
        _job_store[job_id]["status"] = JobStatus.PROCESSING
        images = []
        for path in file_paths:
            img = Image.open(path)
            if img.mode == 'RGBA':
                img = img.convert('RGB')
            images.append(img.copy())
        out_path = os.path.join(tmp_dir, "images_output.pdf")
        buf = io.BytesIO()
        if len(images) == 1:
            images[0].save(buf, 'PDF', resolution=150.0)
        else:
            images[0].save(buf, 'PDF', resolution=150.0, save_all=True, append_images=images[1:])
        buf.seek(0)
        with open(out_path, 'wb') as f:
            f.write(buf.read())
        _job_store[job_id]["status"] = JobStatus.DONE
        _job_store[job_id]["result_path"] = out_path
    except Exception as e:
        _job_store[job_id]["status"] = JobStatus.ERROR
        _job_store[job_id]["error"] = str(e)


def _run_split_pdf_sync(job_id: str, in_path: str, tmp_dir: str, stem: str, ranges_str: str, merge_ranges: bool = False, max_size_kb: int = 0):
    """Split PDF into files by semicolon-separated page range groups, size limit, or merge ranges."""
    try:
        _job_store[job_id]["status"] = JobStatus.PROCESSING
        parts = []

        if max_size_kb > 0 and ranges_str == "1":
            ranges_str = "all"

        try:
            import fitz
            src_doc = fitz.open(in_path)
            total = len(src_doc)
            groups = [g.strip() for g in ranges_str.split(';') if g.strip()]
            if not groups:
                groups = [ranges_str]

            if max_size_kb > 0:
                max_bytes = max_size_kb * 1024
                flat_pages = []
                for group in groups:
                    flat_pages.extend(_parse_page_ranges(group, total))

                if not flat_pages:
                    raise ValueError("No valid pages found in split range.")

                current_doc = fitz.open()
                part_idx = 1

                for p in flat_pages:
                    current_doc.insert_pdf(src_doc, from_page=p, to_page=p)
                    size = len(current_doc.tobytes())
                    if current_doc.page_count > 1 and size > max_bytes:
                        current_doc.delete_page(-1)
                        part_filename = f"{stem}_part{part_idx}.pdf"
                        part_path = os.path.join(tmp_dir, part_filename)
                        current_doc.save(part_path)
                        current_doc.close()
                        parts.append((part_path, part_filename))
                        part_idx += 1

                        current_doc = fitz.open()
                        current_doc.insert_pdf(src_doc, from_page=p, to_page=p)

                if current_doc.page_count > 0:
                    part_filename = f"{stem}_part{part_idx}.pdf"
                    part_path = os.path.join(tmp_dir, part_filename)
                    current_doc.save(part_path)
                    current_doc.close()
                    parts.append((part_path, part_filename))
                else:
                    current_doc.close()

            elif merge_ranges:
                merged_doc = fitz.open()
                for group in groups:
                    pages = _parse_page_ranges(group, total)
                    for p in pages:
                        merged_doc.insert_pdf(src_doc, from_page=p, to_page=p)

                if merged_doc.page_count > 0:
                    part_filename = f"{stem}_merged.pdf"
                    part_path = os.path.join(tmp_dir, part_filename)
                    merged_doc.save(part_path)
                    merged_doc.close()
                    parts.append((part_path, part_filename))
                else:
                    merged_doc.close()

            else:
                for idx, group in enumerate(groups, 1):
                    pages = _parse_page_ranges(group, total)
                    if not pages:
                        continue
                    part_doc = fitz.open()
                    for p in pages:
                        part_doc.insert_pdf(src_doc, from_page=p, to_page=p)
                    part_filename = f"{stem}_part{idx}.pdf"
                    part_path = os.path.join(tmp_dir, part_filename)
                    part_doc.save(part_path)
                    part_doc.close()
                    parts.append((part_path, part_filename))

            src_doc.close()
        except Exception:
            parts = []
            fh = open(in_path, 'rb')
            try:
                reader = PdfReader(fh)
                total = len(reader.pages)
                groups = [g.strip() for g in ranges_str.split(';') if g.strip()]
                if not groups:
                    groups = [ranges_str]

                if max_size_kb > 0:
                    max_bytes = max_size_kb * 1024
                    flat_pages = []
                    for group in groups:
                        flat_pages.extend(_parse_page_ranges(group, total))

                    if not flat_pages:
                        raise ValueError("No valid pages found in split range.")

                    writer = PdfWriter()
                    part_idx = 1

                    for p in flat_pages:
                        writer.add_page(reader.pages[p])
                        buf = io.BytesIO()
                        writer.write(buf)
                        size = buf.tell()
                        if len(writer.pages) > 1 and size > max_bytes:
                            prev_writer = PdfWriter()
                            for pg in writer.pages[:-1]:
                                prev_writer.add_page(pg)
                            part_filename = f"{stem}_part{part_idx}.pdf"
                            part_path = os.path.join(tmp_dir, part_filename)
                            with open(part_path, 'wb') as pf:
                                prev_writer.write(pf)
                            parts.append((part_path, part_filename))
                            part_idx += 1

                            writer = PdfWriter()
                            writer.add_page(reader.pages[p])

                    if len(writer.pages) > 0:
                        part_filename = f"{stem}_part{part_idx}.pdf"
                        part_path = os.path.join(tmp_dir, part_filename)
                        with open(part_path, 'wb') as pf:
                            writer.write(pf)
                        parts.append((part_path, part_filename))

                elif merge_ranges:
                    writer = PdfWriter()
                    for group in groups:
                        pages = _parse_page_ranges(group, total)
                        for p in pages:
                            writer.add_page(reader.pages[p])

                    if len(writer.pages) > 0:
                        part_filename = f"{stem}_merged.pdf"
                        part_path = os.path.join(tmp_dir, part_filename)
                        with open(part_path, 'wb') as pf:
                            writer.write(pf)
                        parts.append((part_path, part_filename))

                else:
                    for idx, group in enumerate(groups, 1):
                        pages = _parse_page_ranges(group, total)
                        if not pages:
                            continue
                        writer = PdfWriter()
                        for p in pages:
                            writer.add_page(reader.pages[p])

                        part_filename = f"{stem}_part{idx}.pdf"
                        part_path = os.path.join(tmp_dir, part_filename)
                        with open(part_path, 'wb') as pf:
                            writer.write(pf)
                        parts.append((part_path, part_filename))
            finally:
                fh.close()

        if not parts:
            raise ValueError("No valid pages found in split range.")

        if len(parts) == 1:
            _job_store[job_id]["result_path"] = parts[0][0]
            _job_store[job_id]["filename"] = parts[0][1]
            _job_store[job_id]["content_type"] = "application/pdf"
        else:
            clean_stem = Path(stem).name
            zip_path = os.path.join(tmp_dir, f"{clean_stem}_split.zip")
            used_names = set()
            with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zf:
                for part_path, part_filename in parts:
                    clean_part_filename = Path(part_filename).name
                    stem_part = Path(clean_part_filename).stem
                    ext_part = Path(clean_part_filename).suffix
                    final_name = clean_part_filename
                    counter = 1
                    while final_name in used_names:
                        final_name = f"{stem_part}_{counter}{ext_part}"
                        counter += 1
                    used_names.add(final_name)
                    zf.write(part_path, arcname=final_name)
            _job_store[job_id]["result_path"] = zip_path
            _job_store[job_id]["filename"] = f"{clean_stem}_split.zip"
            _job_store[job_id]["content_type"] = "application/zip"

        _job_store[job_id]["status"] = JobStatus.DONE
    except Exception as e:
        _job_store[job_id]["status"] = JobStatus.ERROR
        _job_store[job_id]["error"] = str(e)


def _run_remove_pages_sync(job_id: str, in_path: str, tmp_dir: str, stem: str, pages_str: str):
    """Remove specified pages from PDF."""
    try:
        _job_store[job_id]["status"] = JobStatus.PROCESSING
        out_path = os.path.join(tmp_dir, f"{stem}_removed.pdf")
        try:
            import fitz
            doc = fitz.open(in_path)
            total = len(doc)
            to_remove = set(_parse_page_ranges(pages_str, total))
            to_keep = [i for i in range(total) if i not in to_remove]
            doc.select(to_keep)
            doc.save(out_path)
            doc.close()
        except Exception:
            fh = open(in_path, 'rb')
            try:
                reader = PdfReader(fh)
                total = len(reader.pages)
                to_remove = set(_parse_page_ranges(pages_str, total))
                writer = PdfWriter()
                for i, page in enumerate(reader.pages):
                    if i not in to_remove:
                        writer.add_page(page)
                with open(out_path, 'wb') as f:
                    writer.write(f)
            finally:
                fh.close()
        _job_store[job_id]["status"] = JobStatus.DONE
        _job_store[job_id]["result_path"] = out_path
    except Exception as e:
        _job_store[job_id]["status"] = JobStatus.ERROR
        _job_store[job_id]["error"] = str(e)


def _run_extract_pages_sync(job_id: str, in_path: str, tmp_dir: str, stem: str, pages_str: str):
    """Extract (keep) specified pages from PDF."""
    try:
        _job_store[job_id]["status"] = JobStatus.PROCESSING
        out_path = os.path.join(tmp_dir, f"{stem}_extracted.pdf")
        try:
            import fitz
            doc = fitz.open(in_path)
            total = len(doc)
            to_keep = _parse_page_ranges(pages_str, total)
            doc.select(to_keep)
            doc.save(out_path)
            doc.close()
        except Exception:
            fh = open(in_path, 'rb')
            try:
                reader = PdfReader(fh)
                total = len(reader.pages)
                to_keep = _parse_page_ranges(pages_str, total)
                writer = PdfWriter()
                for i in to_keep:
                    writer.add_page(reader.pages[i])
                with open(out_path, 'wb') as f:
                    writer.write(f)
            finally:
                fh.close()
        _job_store[job_id]["status"] = JobStatus.DONE
        _job_store[job_id]["result_path"] = out_path
    except Exception as e:
        _job_store[job_id]["status"] = JobStatus.ERROR
        _job_store[job_id]["error"] = str(e)


def _run_organize_pdf_sync(
    job_id: str,
    file_paths: Union[str, List[str]],
    tmp_dir: str,
    stem: str,
    order_str: str,
    add_page_numbers: bool = False,
    page_number_position: str = "bottom-center",
    page_number_format: str = "Page {page} of {total}",
    start_number: int = 1,
):
    """Reorder and assemble PDF pages from one or multiple PDFs to a custom order, optionally adding page numbers."""
    try:
        _job_store[job_id]["status"] = JobStatus.PROCESSING
        out_path = os.path.join(tmp_dir, f"{stem}_organized.pdf")
        if isinstance(file_paths, str):
            paths = [file_paths]
        else:
            paths = file_paths

        try:
            import fitz
            docs = [fitz.open(fp) for fp in paths]
            doc_page_counts = [len(d) for d in docs]
            order_tuples = _parse_multi_pdf_order(order_str, doc_page_counts)

            out_doc = fitz.open()
            for f_idx, p_idx in order_tuples:
                out_doc.insert_pdf(docs[f_idx], from_page=p_idx, to_page=p_idx)

            total_output_pages = len(out_doc)
            if add_page_numbers and total_output_pages > 0:
                pos = (page_number_position or "bottom-center").lower().replace("_", "-").strip()
                for i in range(total_output_pages):
                    page = out_doc[i]
                    page_num = start_number + i
                    try:
                        text = page_number_format.format(page=page_num, total=total_output_pages)
                    except Exception:
                        text = f"Page {page_num} of {total_output_pages}"

                    rect = page.rect
                    margin_x = 36
                    margin_y = 20
                    box_h = 20

                    if "top" in pos:
                        y0 = margin_y
                        y1 = margin_y + box_h
                    else:  # bottom
                        y0 = rect.height - margin_y - box_h
                        y1 = rect.height - margin_y

                    if "left" in pos:
                        box = fitz.Rect(margin_x, y0, margin_x + 200, y1)
                        align = fitz.TEXT_ALIGN_LEFT
                    elif "right" in pos:
                        box = fitz.Rect(rect.width - margin_x - 200, y0, rect.width - margin_x, y1)
                        align = fitz.TEXT_ALIGN_RIGHT
                    else:  # center / middle
                        box = fitz.Rect(margin_x, y0, rect.width - margin_x, y1)
                        align = fitz.TEXT_ALIGN_CENTER

                    page.insert_textbox(box, text, fontsize=10, align=align)

            out_doc.save(out_path)
            out_doc.close()
            for d in docs:
                d.close()
        except Exception:
            from PyPDF2 import PdfReader, PdfWriter
            handles = [open(fp, 'rb') for fp in paths]
            try:
                readers = [PdfReader(h) for h in handles]
                doc_page_counts = [len(r.pages) for r in readers]
                order_tuples = _parse_multi_pdf_order(order_str, doc_page_counts)
                writer = PdfWriter()
                for f_idx, p_idx in order_tuples:
                    writer.add_page(readers[f_idx].pages[p_idx])
                with open(out_path, 'wb') as f:
                    writer.write(f)
            finally:
                for h in handles:
                    h.close()

        _job_store[job_id]["status"] = JobStatus.DONE
        _job_store[job_id]["result_path"] = out_path
    except Exception as e:
        _job_store[job_id]["status"] = JobStatus.ERROR
        _job_store[job_id]["error"] = str(e)


def _run_rotate_pdf_sync(job_id: str, in_path: str, tmp_dir: str, stem: str, rotation: int, pages_str: str, rotations_json: Optional[str] = None):
    """Rotate selected (or all) pages of a PDF, or apply per-page rotation angles."""
    try:
        _job_store[job_id]["status"] = JobStatus.PROCESSING
        out_path = os.path.join(tmp_dir, f"{stem}_rotated.pdf")

        rot_map = {}
        if rotations_json:
            try:
                parsed = json.loads(rotations_json)
                if isinstance(parsed, dict):
                    for k, v in parsed.items():
                        rot_map[int(k)] = int(v) % 360
            except Exception as pe:
                print(f"Warning: Failed to parse rotations_json: {pe}")

        try:
            import fitz
            doc = fitz.open(in_path)
            total = len(doc)

            if rot_map:
                for i in range(total):
                    p_num = i + 1
                    if p_num in rot_map and rot_map[p_num] != 0:
                        angle = rot_map[p_num]
                        page = doc[i]
                        page.set_rotation((page.rotation + angle) % 360)
            else:
                if pages_str.strip().lower() in ('all', ''):
                    to_rotate = set(range(total))
                else:
                    to_rotate = set(_parse_page_ranges(pages_str, total))
                for i in to_rotate:
                    page = doc[i]
                    page.set_rotation((page.rotation + rotation) % 360)
            doc.save(out_path)
            doc.close()
        except Exception:
            fh = open(in_path, 'rb')
            try:
                reader = PdfReader(fh)
                total = len(reader.pages)
                writer = PdfWriter()

                if rot_map:
                    for i, page in enumerate(reader.pages):
                        p_num = i + 1
                        if p_num in rot_map and rot_map[p_num] != 0:
                            angle = rot_map[p_num]
                            page.rotate(angle)
                        writer.add_page(page)
                else:
                    if pages_str.strip().lower() in ('all', ''):
                        to_rotate = set(range(total))
                    else:
                        to_rotate = set(_parse_page_ranges(pages_str, total))
                    for i, page in enumerate(reader.pages):
                        if i in to_rotate:
                            page.rotate(rotation)
                        writer.add_page(page)
                with open(out_path, 'wb') as f:
                    writer.write(f)
            finally:
                fh.close()
        _job_store[job_id]["status"] = JobStatus.DONE
        _job_store[job_id]["result_path"] = out_path
    except Exception as e:
        _job_store[job_id]["status"] = JobStatus.ERROR
        _job_store[job_id]["error"] = str(e)


def _hex_to_rgb(hex_str: str) -> tuple:
    hex_str = hex_str.lstrip('#').strip()
    if len(hex_str) == 3:
        hex_str = ''.join([c * 2 for c in hex_str])
    if len(hex_str) == 6:
        try:
            r = int(hex_str[0:2], 16) / 255.0
            g = int(hex_str[2:4], 16) / 255.0
            b = int(hex_str[4:6], 16) / 255.0
            return (r, g, b)
        except ValueError:
            pass
    return (1.0, 0.0, 0.0)


def _get_watermark_font_name(font_family: str, bold: bool, italic: bool) -> str:
    family = font_family.strip().lower()
    if "times" in family:
        if bold and italic:
            return "Times-BoldItalic"
        elif bold:
            return "Times-Bold"
        elif italic:
            return "Times-Italic"
        return "Times-Roman"
    elif "courier" in family:
        if bold and italic:
            return "Courier-BoldOblique"
        elif bold:
            return "Courier-Bold"
        elif italic:
            return "Courier-Oblique"
        return "Courier"
    else:
        if bold and italic:
            return "Helvetica-BoldOblique"
        elif bold:
            return "Helvetica-Bold"
        elif italic:
            return "Helvetica-Oblique"
        return "Helvetica"


def _run_add_watermark_sync(
    job_id: str,
    in_path: str,
    tmp_dir: str,
    stem: str,
    watermark_type: str = "text",
    text: str = "CONFIDENTIAL",
    image_path: Optional[str] = None,
    font_family: str = "Helvetica",
    font_size: int = 36,
    bold: bool = False,
    italic: bool = False,
    color: str = "#FF0000",
    position: str = "center",
    opacity: float = 0.5,
    rotation: int = 45,
    from_page: int = 1,
    to_page: int = 0,
    layer: str = "over",
    scale: float = 1.0,
):
    """Add text or image watermark to specified pages using PyMuPDF."""
    try:
        import fitz
        _job_store[job_id]["status"] = JobStatus.PROCESSING
        doc = fitz.open(in_path)
        total_pages = len(doc)

        start_p = max(0, from_page - 1)
        if to_page <= 0 or to_page > total_pages:
            end_p = total_pages
        else:
            end_p = min(to_page, total_pages)

        overlay_flag = (layer == "over")
        clamped_scale = max(0.2, min(5.0, scale))

        if watermark_type == "text":
            font_size = max(1, int(font_size * clamped_scale))
            rgb = _hex_to_rgb(color)
            font_name = _get_watermark_font_name(font_family, bold, italic)
            try:
                font_obj = fitz.Font(font_name)
            except Exception:
                font_obj = fitz.Font("helv")

            for page_num in range(start_p, end_p):
                if page_num >= total_pages:
                    break
                page = doc[page_num]
                pw, ph = page.rect.width, page.rect.height
                tw = font_obj.text_length(text, font_size)
                th = float(font_size)
                margin = 50.0

                pos_lower = position.lower()
                if pos_lower == "mosaic":
                    xs = [pw / 6.0 - tw / 2.0, pw / 2.0 - tw / 2.0, 5.0 * pw / 6.0 - tw / 2.0]
                    ys = [ph / 6.0 + th / 3.0, ph / 2.0 + th / 3.0, 5.0 * ph / 6.0 + th / 3.0]
                    pts = [(x, y) for x in xs for y in ys]
                else:
                    x_left = margin
                    x_center = max(0.0, (pw - tw) / 2.0)
                    x_right = max(0.0, pw - tw - margin)

                    y_top = margin + th
                    y_center = ph / 2.0 + th / 3.0
                    y_bottom = max(margin, ph - margin)

                    coords_map = {
                        "top-left": (x_left, y_top),
                        "top-center": (x_center, y_top),
                        "top-right": (x_right, y_top),
                        "center-left": (x_left, y_center),
                        "center": (x_center, y_center),
                        "center-right": (x_right, y_center),
                        "bottom-left": (x_left, y_bottom),
                        "bottom-center": (x_center, y_bottom),
                        "bottom-right": (x_right, y_bottom),
                    }
                    pt = coords_map.get(pos_lower, (x_center, y_center))
                    pts = [pt]

                for pt in pts:
                    if rotation % 360 == 0:
                        page.insert_text(
                            pt, text,
                            fontname=font_name,
                            fontsize=font_size,
                            color=rgb,
                            fill_opacity=opacity,
                            overlay=overlay_flag,
                        )
                    else:
                        cx = pt[0] + tw / 2.0
                        cy = pt[1] - th / 3.0
                        morph_tuple = (fitz.Point(cx, cy), fitz.Matrix(-rotation))
                        page.insert_text(
                            pt, text,
                            fontname=font_name,
                            fontsize=font_size,
                            color=rgb,
                            fill_opacity=opacity,
                            morph=morph_tuple,
                            overlay=overlay_flag,
                        )

        elif watermark_type == "image":
            if not image_path or not os.path.exists(image_path):
                raise ValueError("Image file not found for image watermark.")

            from PIL import Image
            img = Image.open(image_path)
            if img.mode != "RGBA":
                img = img.convert("RGBA")

            if rotation % 360 != 0:
                img = img.rotate(-rotation, expand=True, resample=Image.BICUBIC)

            if opacity < 1.0:
                r, g, b, a = img.split()
                a = a.point(lambda p: int(p * max(0.0, min(1.0, opacity))))
                img = Image.merge("RGBA", (r, g, b, a))

            img_buf = io.BytesIO()
            img.save(img_buf, format="PNG")
            img_bytes = img_buf.getvalue()
            iw, ih = img.size

            for page_num in range(start_p, end_p):
                if page_num >= total_pages:
                    break
                page = doc[page_num]
                pw, ph = page.rect.width, page.rect.height

                pos_lower = position.lower()
                if pos_lower == "mosaic":
                    max_w, max_h = pw * 0.25 * clamped_scale, ph * 0.25 * clamped_scale
                else:
                    max_w, max_h = pw * 0.4 * clamped_scale, ph * 0.4 * clamped_scale

                img_scale = min(max_w / max(1, iw), max_h / max(1, ih), 1.0)
                dw = iw * img_scale
                dh = ih * img_scale
                margin = 50.0

                if pos_lower == "mosaic":
                    xs = [pw / 6.0, pw / 2.0, 5.0 * pw / 6.0]
                    ys = [ph / 6.0, ph / 2.0, 5.0 * ph / 6.0]
                    center_pts = [(x, y) for x in xs for y in ys]
                else:
                    cx_left = margin + dw / 2.0
                    cx_center = pw / 2.0
                    cx_right = pw - margin - dw / 2.0

                    cy_top = margin + dh / 2.0
                    cy_center = ph / 2.0
                    cy_bottom = ph - margin - dh / 2.0

                    centers_map = {
                        "top-left": (cx_left, cy_top),
                        "top-center": (cx_center, cy_top),
                        "top-right": (cx_right, cy_top),
                        "center-left": (cx_left, cy_center),
                        "center": (cx_center, cy_center),
                        "center-right": (cx_right, cy_center),
                        "bottom-left": (cx_left, cy_bottom),
                        "bottom-center": (cx_center, cy_bottom),
                        "bottom-right": (cx_right, cy_bottom),
                    }
                    c_pt = centers_map.get(pos_lower, (cx_center, cy_center))
                    center_pts = [c_pt]

                for cx, cy in center_pts:
                    rect = fitz.Rect(cx - dw / 2.0, cy - dh / 2.0, cx + dw / 2.0, cy + dh / 2.0)
                    page.insert_image(rect, stream=img_bytes, overlay=overlay_flag)

        out_path = os.path.join(tmp_dir, f"{stem}_watermarked.pdf")
        doc.save(out_path, garbage=3, deflate=True)
        doc.close()
        _job_store[job_id]["status"] = JobStatus.DONE
        _job_store[job_id]["result_path"] = out_path
    except Exception as e:
        _job_store[job_id]["status"] = JobStatus.ERROR
        _job_store[job_id]["error"] = str(e)


def _run_add_page_numbers_sync(job_id: str, in_path: str, tmp_dir: str, stem: str, position: str, start: int):
    """Add page numbers to each page using PyMuPDF."""
    try:
        import fitz
        _job_store[job_id]["status"] = JobStatus.PROCESSING
        doc = fitz.open(in_path)
        for i, page in enumerate(doc):
            num = i + start
            w, h = page.rect.width, page.rect.height
            pt = fitz.Point(w / 2 - 10, 20 if position == 'top' else h - 15)
            page.insert_text(pt, str(num), fontsize=11, color=(0, 0, 0))
        out_path = os.path.join(tmp_dir, f"{stem}_numbered.pdf")
        doc.save(out_path, garbage=3, deflate=True)
        doc.close()
        _job_store[job_id]["status"] = JobStatus.DONE
        _job_store[job_id]["result_path"] = out_path
    except Exception as e:
        _job_store[job_id]["status"] = JobStatus.ERROR
        _job_store[job_id]["error"] = str(e)





def _run_repair_pdf_sync(job_id: str, in_path: str, tmp_dir: str, stem: str):
    """Repair PDF using PyMuPDF / pikepdf — runs in thread pool."""
    try:
        _job_store[job_id]["status"] = JobStatus.PROCESSING
        out_path = os.path.join(tmp_dir, f"{stem}_repaired.pdf")
        repaired = False

        try:
            import fitz
            doc = fitz.open(in_path)
            doc.save(out_path, garbage=3, deflate=True)
            doc.close()
            repaired = os.path.isfile(out_path) and os.path.getsize(out_path) > 0
        except Exception:
            repaired = False

        if not repaired and PIKEPDF_AVAILABLE and pikepdf:
            try:
                with pikepdf.Pdf.open(in_path) as pdf:
                    pdf.save(out_path)
                repaired = os.path.isfile(out_path) and os.path.getsize(out_path) > 0
            except Exception:
                repaired = False

        if not repaired:
            _job_store[job_id]["status"] = JobStatus.ERROR
            _job_store[job_id]["error"] = "PDF file is too severely corrupt to repair."
            return

        _job_store[job_id]["status"] = JobStatus.DONE
        _job_store[job_id]["result_path"] = out_path
    except Exception as e:
        _job_store[job_id]["status"] = JobStatus.ERROR
        _job_store[job_id]["error"] = str(e)


def _run_pdf_to_pptx_sync(job_id: str, in_path: str, tmp_dir: str, stem: str):
    """Convert PDF to PowerPoint (.pptx) with editable text boxes and native pictures, with fallback to raster images for scanned pages."""
    try:
        import fitz
        import re
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor

        _job_store[job_id]["status"] = JobStatus.PROCESSING
        doc = fitz.open(in_path)
        prs = Presentation()

        if len(doc) > 0:
            first_rect = doc[0].rect
            prs.slide_width = Inches(first_rect.width / 72.0)
            prs.slide_height = Inches(first_rect.height / 72.0)
        else:
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

        blank_slide_layout = prs.slide_layouts[6]
        EMU_PER_PT = 12700

        def _map_font(font_name: str) -> dict:
            fn_lower = font_name.lower()
            is_bold = any(k in fn_lower for k in ["bold", "black", "heavy", "demi"])
            is_italic = any(k in fn_lower for k in ["italic", "oblique", "slant"])

            clean = font_name
            for s in ["-Bold", "-Italic", "-BoldItalic", "-Oblique", "MT", "PS", "Regular", "Roman"]:
                clean = re.sub(re.escape(s) + r"$", "", clean, flags=re.I)
                clean = clean.replace(s, "")
            clean = clean.strip("- ").lower()

            font_map = {
                "arial": "Arial",
                "helvetica": "Arial",
                "times": "Times New Roman",
                "timesnewroman": "Times New Roman",
                "courier": "Courier New",
                "couriernew": "Courier New",
                "calibri": "Calibri",
                "verdana": "Verdana",
                "georgia": "Georgia",
                "trebuchet": "Trebuchet MS",
                "garamond": "Garamond",
            }
            return {
                "name": font_map.get(clean, "Calibri"),
                "bold": is_bold,
                "italic": is_italic,
            }

        img_dir = os.path.join(tmp_dir, "extracted_imgs")
        os.makedirs(img_dir, exist_ok=True)

        for page in doc:
            slide = prs.slides.add_slide(blank_slide_layout)
            page_w = page.rect.width
            page_h = page.rect.height

            # Extract structured text dictionary
            page_dict = page.get_text("dict")
            blocks = page_dict.get("blocks", [])

            # Calculate total text length on page to check if scanned
            text_blocks = [b for b in blocks if b.get("type") == 0]
            total_text_len = sum(
                sum(len(s.get("text", "").strip()) for l in b.get("lines", []) for s in l.get("spans", []))
                for b in text_blocks
            )

            # Scanned / Image-only page fallback
            if total_text_len == 0:
                pix = page.get_pixmap(dpi=150)
                img_path = os.path.join(tmp_dir, f"page_{page.number}.png")
                pix.save(img_path)
                slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
                continue

            # 1. Embedded Images
            try:
                seen_xrefs = set()
                for img_info in page.get_images(full=True):
                    xref = img_info[0]
                    if xref in seen_xrefs:
                        continue
                    seen_xrefs.add(xref)

                    rects = page.get_image_rects(xref)
                    if not rects:
                        continue

                    base_img = doc.extract_image(xref)
                    if not base_img or not base_img.get("image"):
                        continue

                    img_ext = base_img.get("ext", "png")
                    img_filename = os.path.join(img_dir, f"p{page.number}_xref{xref}.{img_ext}")
                    with open(img_filename, "wb") as f:
                        f.write(base_img["image"])

                    for r in rects:
                        l_emu = int(r.x0 * EMU_PER_PT)
                        t_emu = int(r.y0 * EMU_PER_PT)
                        w_emu = int(r.width * EMU_PER_PT)
                        h_emu = int(r.height * EMU_PER_PT)
                        if w_emu > 0 and h_emu > 0:
                            slide.shapes.add_picture(img_filename, Emu(l_emu), Emu(t_emu), Emu(w_emu), Emu(h_emu))
            except Exception:
                pass

            # 2. Text Boxes (Line-level grouped text frames)
            for b in text_blocks:
                bbox = b.get("bbox", (0, 0, page_w, page_h))
                for line in b.get("lines", []):
                    spans = line.get("spans", [])
                    if not spans:
                        continue

                    line_bbox = line.get("bbox", bbox)
                    lx0, ly0, lx1, ly1 = line_bbox
                    
                    lw = max(lx1 - lx0, 10.0)
                    lh = max(ly1 - ly0, 10.0)

                    tx_box = slide.shapes.add_textbox(
                        Emu(int(lx0 * EMU_PER_PT)),
                        Emu(int(ly0 * EMU_PER_PT)),
                        Emu(int(lw * EMU_PER_PT)),
                        Emu(int(lh * EMU_PER_PT))
                    )

                    tf = tx_box.text_frame
                    tf.word_wrap = True
                    tf.margin_left = Emu(0)
                    tf.margin_right = Emu(0)
                    tf.margin_top = Emu(0)
                    tf.margin_bottom = Emu(0)

                    p = tf.paragraphs[0]
                    p.space_before = Pt(0)
                    p.space_after = Pt(0)

                    for idx, span in enumerate(spans):
                        stext = span.get("text", "")
                        if not stext:
                            continue

                        run = p.add_run()
                        run.text = stext

                        # Font Styling
                        fm = _map_font(span.get("font", "Calibri"))
                        run.font.name = fm["name"]
                        run.font.size = Pt(span.get("size", 11))
                        run.font.bold = fm["bold"]
                        run.font.italic = fm["italic"]

                        color_int = span.get("color", 0)
                        r_val = (color_int >> 16) & 0xFF
                        g_val = (color_int >> 8) & 0xFF
                        b_val = color_int & 0xFF
                        run.font.color.rgb = RGBColor(r_val, g_val, b_val)

                    tx_box.fill.background()
                    tx_box.line.fill.background()

        doc.close()
        out_path = os.path.join(tmp_dir, f"{stem}.pptx")
        prs.save(out_path)

        _job_store[job_id]["status"] = JobStatus.DONE
        _job_store[job_id]["result_path"] = out_path
    except Exception as e:
        _job_store[job_id]["status"] = JobStatus.ERROR
        _job_store[job_id]["error"] = str(e)


def _extract_spatial_table_from_fitz_page(page) -> list:
    """Extract multi-column borderless tables using word bounding box spatial clustering."""
    words = page.get_text("words")
    if not words:
        return []

    lines = []
    for w in sorted(words, key=lambda item: (item[1], item[0])):
        x0, y0, x1, y1, text = w[0], w[1], w[2], w[3], w[4]
        matched = False
        for line in lines:
            if abs(line['y'] - y0) < 4.5:
                line['words'].append(w)
                matched = True
                break
        if not matched:
            lines.append({'y': y0, 'words': [w]})

    all_gaps = []
    for line in lines:
        l_words = sorted(line['words'], key=lambda item: item[0])
        for i in range(len(l_words) - 1):
            gap_start = l_words[i][2]
            gap_end = l_words[i+1][0]
            if gap_end - gap_start > 14:
                all_gaps.append((gap_start, gap_end))

    col_breaks = []
    if all_gaps:
        all_gaps.sort(key=lambda g: g[0])
        merged_gaps = []
        for g in all_gaps:
            if not merged_gaps:
                merged_gaps.append(list(g))
            else:
                last = merged_gaps[-1]
                if g[0] <= last[1] + 15:
                    last[1] = max(last[1], g[1])
                    last[0] = min(last[0], g[0])
                else:
                    merged_gaps.append(list(g))
        col_breaks = [(g[0] + g[1]) / 2 for g in merged_gaps if (g[1] - g[0]) >= 12]

    table_rows = []
    for line in lines:
        l_words = sorted(line['words'], key=lambda item: item[0])
        num_cols = len(col_breaks) + 1
        row = [""] * num_cols
        for w in l_words:
            wx0 = w[0]
            col_idx = 0
            for b_idx, b_x in enumerate(col_breaks):
                if wx0 >= b_x:
                    col_idx = b_idx + 1
            if row[col_idx]:
                row[col_idx] += " " + w[4]
            else:
                row[col_idx] = w[4]
        table_rows.append(row)

    return table_rows


def _run_pdf_to_excel_sync(job_id: str, in_path: str, tmp_dir: str, stem: str):
    """Extract tables from PDF to Excel (.xlsx) using PyMuPDF (fitz) and openpyxl, with spatial column clustering."""
    try:
        _job_store[job_id]["status"] = JobStatus.PROCESSING
        import openpyxl
        wb = openpyxl.Workbook()
        wb.remove(wb.active)

        tables_found = 0

        # Primary extraction engine: PyMuPDF (fitz) vector table finder
        try:
            import fitz
            doc = fitz.open(in_path)
            for i, page in enumerate(doc, 1):
                try:
                    tabs = page.find_tables()
                    if tabs and tabs.tables:
                        for t_idx, tab in enumerate(tabs.tables, 1):
                            data = tab.extract()
                            if not data:
                                continue
                            tables_found += 1
                            sheet_name = f"Page{i}_Table{t_idx}"[:31]
                            ws = wb.create_sheet(title=sheet_name)
                            for row in data:
                                ws.append([cell if cell is not None else "" for cell in row])
                except Exception:
                    pass
            doc.close()
        except Exception:
            pass

        # Secondary extraction engine: Spatial layout column clusterer (for borderless tables like rosters)
        if tables_found == 0:
            try:
                import fitz
                doc = fitz.open(in_path)
                for i, page in enumerate(doc, 1):
                    rows = _extract_spatial_table_from_fitz_page(page)
                    if rows:
                        tables_found += 1
                        sheet_name = f"Page_{i}"[:31]
                        ws = wb.create_sheet(title=sheet_name)
                        for r in rows:
                            ws.append(r)
                doc.close()
            except Exception:
                pass

        # Fallback extraction engine: pdfplumber if installed
        if tables_found == 0:
            try:
                import pdfplumber
                with pdfplumber.open(in_path) as pdf:
                    for i, page in enumerate(pdf.pages, 1):
                        tables = page.extract_tables()
                        for t_idx, table in enumerate(tables, 1):
                            if not table:
                                continue
                            tables_found += 1
                            sheet_name = f"Page{i}_Table{t_idx}"[:31]
                            ws = wb.create_sheet(title=sheet_name)
                            for row in table:
                                ws.append([cell if cell is not None else "" for cell in row])
            except Exception:
                pass

        # Fallback text/line extraction if no structured tables were detected
        if tables_found == 0:
            try:
                import fitz
                import re
                doc = fitz.open(in_path)
                for i, page in enumerate(doc, 1):
                    sheet_name = f"Page_{i}"[:31]
                    ws = wb.create_sheet(title=sheet_name)
                    text = page.get_text("text") or ""
                    for line in text.splitlines():
                        if line.strip():
                            cols = re.split(r'\t+|\s{2,}', line.strip())
                            ws.append(cols)
                doc.close()
            except Exception:
                pass

        if not wb.sheetnames:
            ws = wb.create_sheet(title="Sheet1")
            ws.append(["No text or tables found in PDF"])

        out_path = os.path.join(tmp_dir, f"{stem}.xlsx")
        wb.save(out_path)

        _job_store[job_id]["status"] = JobStatus.DONE
        _job_store[job_id]["result_path"] = out_path
    except Exception as e:
        _job_store[job_id]["status"] = JobStatus.ERROR
        _job_store[job_id]["error"] = str(e)


def _run_office_to_pdf_sync(job_id: str, in_path: str, tmp_dir: str, soffice: str):
    """Convert PPTX/Excel to PDF via LibreOffice in thread pool."""
    try:
        _job_store[job_id]["status"] = JobStatus.PROCESSING
        result = subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf", "--outdir", tmp_dir, in_path],
            capture_output=True,
            text=True,
        )
        if result.returncode != 0:
            _job_store[job_id]["status"] = JobStatus.ERROR
            _job_store[job_id]["error"] = f"LibreOffice conversion failed: {result.stderr}"
            return

        from pathlib import Path
        pdf_path = os.path.join(tmp_dir, Path(in_path).stem + ".pdf")
        if not os.path.isfile(pdf_path):
            pdf_files = [f for f in os.listdir(tmp_dir) if f.endswith(".pdf")]
            if pdf_files:
                pdf_path = os.path.join(tmp_dir, pdf_files[0])
            else:
                _job_store[job_id]["status"] = JobStatus.ERROR
                _job_store[job_id]["error"] = "LibreOffice did not produce a PDF output."
                return

        _job_store[job_id]["status"] = JobStatus.DONE
        _job_store[job_id]["result_path"] = pdf_path
    except Exception as e:
        _job_store[job_id]["status"] = JobStatus.ERROR
        _job_store[job_id]["error"] = str(e)


def _run_flatten_pdf_sync(job_id: str, in_path: str, tmp_dir: str, stem: str):
    """Flatten interactive form fields into static PDF content using PyMuPDF."""
    try:
        import fitz
        _job_store[job_id]["status"] = JobStatus.PROCESSING
        doc = fitz.open(in_path)
        doc.bake()
        out_path = os.path.join(tmp_dir, f"{stem}_flat.pdf")
        doc.save(out_path, garbage=3, deflate=True)
        doc.close()

        _job_store[job_id]["status"] = JobStatus.DONE
        _job_store[job_id]["result_path"] = out_path
    except Exception as e:
        _job_store[job_id]["status"] = JobStatus.ERROR
        _job_store[job_id]["error"] = str(e)


def _is_truthy(val):
    if isinstance(val, bool):
        return val
    if isinstance(val, str):
        return val.strip().lower() in ("true", "1", "yes", "on")
    return bool(val)


def _run_lock_pdf_sync(job_id: str, in_path: str, tmp_dir: str, stem: str, password: str, lock_mode: str, allow_print: bool, allow_copy: bool):
    """Encrypt PDF with 256-bit AES using PyMuPDF.
    
    lock_mode:
      - 'recoverable' / 'permission': Requires password to open & view PDF, BUT owner_pw is empty so it CAN BE UNLOCKED automatically by Unlock PDF without knowing the password!
      - 'open' / 'strict': Requires password to open & view PDF, AND sets secret owner_pw so it CANNOT be unlocked without the password.
    """
    try:
        import fitz
        _job_store[job_id]["status"] = JobStatus.PROCESSING
        doc = fitz.open(in_path)

        # Base permissions mask: -4 allows full functionality
        perm = -4

        is_print_allowed = _is_truthy(allow_print)
        is_copy_allowed = _is_truthy(allow_copy)

        if not is_print_allowed:
            perm &= ~(fitz.PDF_PERM_PRINT | fitz.PDF_PERM_PRINT_HQ)
        if not is_copy_allowed:
            perm &= ~(fitz.PDF_PERM_COPY | fitz.PDF_PERM_ACCESSIBILITY)

        if lock_mode in ("permission", "recoverable"):
            # Recoverable open lock: user must enter password to open in PDF viewers, but owner_pw = "" allows automated Unlock PDF recovery
            user_pw = password
            owner_pw = ""
        else:  # 'open' / 'strict'
            # Strict open lock: user must enter password to open, and owner_pw is secret (cannot be unlocked without password)
            user_pw = password
            owner_pw = password + "_owner"

        out_path = os.path.join(tmp_dir, f"{stem}_protected.pdf")
        doc.save(
            out_path,
            encryption=fitz.PDF_ENCRYPT_AES_256,
            user_pw=user_pw,
            owner_pw=owner_pw,
            permissions=perm,
            deflate=True,
        )
        doc.close()

        _job_store[job_id]["status"] = JobStatus.DONE
        _job_store[job_id]["result_path"] = out_path
    except Exception as e:
        _job_store[job_id]["status"] = JobStatus.ERROR
        _job_store[job_id]["error"] = str(e)


def _run_unlock_pdf_sync(job_id: str, in_path: str, tmp_dir: str, stem: str, password: str):
    """Remove PDF encryption, password protection, and security restrictions."""
    try:
        _job_store[job_id]["status"] = JobStatus.PROCESSING
        out_path = os.path.join(tmp_dir, f"{stem}_unlocked.pdf")
        unlocked = False

        # Method 1: pikepdf (powerful engine for stripping owner/user passwords)
        if PIKEPDF_AVAILABLE and pikepdf:
            try:
                pw_to_try = password.strip() if password else ""
                try:
                    pdf = pikepdf.Pdf.open(in_path, password=pw_to_try)
                except Exception:
                    pdf = None
                    if not pw_to_try:
                        # Try empty password explicitly
                        try:
                            pdf = pikepdf.Pdf.open(in_path, password="")
                        except Exception:
                            pdf = None

                if pdf is not None:
                    pdf.save(out_path)
                    pdf.close()
                    unlocked = os.path.isfile(out_path) and os.path.getsize(out_path) > 0
            except Exception:
                unlocked = False

        # Method 2: PyMuPDF engine fallback
        if not unlocked:
            import fitz
            doc = fitz.open(in_path)
            if doc.is_encrypted:
                auth = doc.authenticate(password or "")
                if not auth and password:
                    # try empty string as fallback for owner passwords
                    auth = doc.authenticate("")
                if not auth:
                    _job_store[job_id]["status"] = JobStatus.ERROR
                    _job_store[job_id]["error"] = "Incorrect password. Enter the correct password to unlock this PDF."
                    doc.close()
                    return
            new_doc = fitz.open()
            new_doc.insert_pdf(doc)
            new_doc.save(out_path, encryption=fitz.PDF_ENCRYPT_NONE, deflate=True)
            new_doc.close()
            doc.close()
            unlocked = os.path.isfile(out_path) and os.path.getsize(out_path) > 0

        if not unlocked:
            _job_store[job_id]["status"] = JobStatus.ERROR
            _job_store[job_id]["error"] = "Incorrect password. Could not unlock the PDF."
            return

        _job_store[job_id]["status"] = JobStatus.DONE
        _job_store[job_id]["result_path"] = out_path
    except Exception as e:
        _job_store[job_id]["status"] = JobStatus.ERROR
        _job_store[job_id]["error"] = str(e)


def _hex_to_fitz_color(hex_color: str):
    """Convert '#RRGGBB' hex string to a normalised (r, g, b) float tuple for PyMuPDF."""
    hex_color = hex_color.lstrip("#")
    if len(hex_color) == 3:
        hex_color = "".join(c * 2 for c in hex_color)
    try:
        r = int(hex_color[0:2], 16) / 255.0
        g = int(hex_color[2:4], 16) / 255.0
        b = int(hex_color[4:6], 16) / 255.0
        return (r, g, b)
    except Exception:
        return (0.0, 0.0, 0.0)


def _run_redact_pdf_sync(job_id: str, in_path: str, tmp_dir: str, stem: str, redactions_json: str):
    """Apply permanent redactions (rectangles + text search) to a PDF using PyMuPDF."""
    import json
    try:
        import fitz
        _job_store[job_id]["status"] = JobStatus.PROCESSING

        try:
            payload = json.loads(redactions_json)
        except Exception:
            _job_store[job_id]["status"] = JobStatus.ERROR
            _job_store[job_id]["error"] = "Invalid redactions JSON payload."
            return

        items = payload.get("items", [])
        doc = fitz.open(in_path)

        for item in items:
            page_idx = int(item.get("page", 0))
            if not (0 <= page_idx < len(doc)):
                continue

            page = doc[page_idx]
            color_hex = item.get("color", "#000000")
            fill_color = _hex_to_fitz_color(color_hex)
            item_type = item.get("type", "rect")

            if item_type == "rect":
                r = item.get("rect", [0, 0, 0, 0])
                if len(r) >= 4:
                    fitz_rect = fitz.Rect(float(r[0]), float(r[1]), float(r[2]), float(r[3]))
                    page.add_redact_annot(fitz_rect, fill=fill_color)

            elif item_type == "text":
                search_text = item.get("text", "").strip()
                if search_text:
                    matches = page.search_for(search_text)
                    for match in matches:
                        page.add_redact_annot(match, fill=fill_color)

        # Apply all redaction annotations — this permanently purges underlying content
        for page in doc:
            page.apply_redactions(images=fitz.PDF_REDACT_IMAGE_NONE)

        out_path = os.path.join(tmp_dir, f"{stem}_redacted.pdf")
        doc.save(out_path, deflate=True)
        doc.close()

        _job_store[job_id]["status"] = JobStatus.DONE
        _job_store[job_id]["result_path"] = out_path
    except Exception as e:
        _job_store[job_id]["status"] = JobStatus.ERROR
        _job_store[job_id]["error"] = str(e)


def _run_html_to_pdf_sync(job_id: str, html_content: str, tmp_dir: str, stem: str):
    """Convert HTML string to PDF using WeasyPrint or PyMuPDF HTML engine."""
    try:
        _job_store[job_id]["status"] = JobStatus.PROCESSING
        out_path = os.path.join(tmp_dir, f"{stem}.pdf")
        converted = False

        if WEASYPRINT_AVAILABLE and weasyprint:
            try:
                weasyprint.HTML(string=html_content).write_pdf(out_path)
                converted = os.path.isfile(out_path) and os.path.getsize(out_path) > 0
            except Exception:
                converted = False

        if not converted:
            try:
                import fitz
                html_bytes = html_content.encode("utf-8") if isinstance(html_content, str) else html_content
                doc = fitz.open("html", html_bytes)
                pdf_bytes = doc.convert_to_pdf()
                doc.close()
                with open(out_path, "wb") as f:
                    f.write(pdf_bytes)
                converted = os.path.isfile(out_path) and os.path.getsize(out_path) > 0
            except Exception:
                converted = False

        if not converted:
            _job_store[job_id]["status"] = JobStatus.ERROR
            _job_store[job_id]["error"] = "HTML to PDF conversion failed."
            return

        _job_store[job_id]["status"] = JobStatus.DONE
        _job_store[job_id]["result_path"] = out_path
    except Exception as e:
        _job_store[job_id]["status"] = JobStatus.ERROR
        _job_store[job_id]["error"] = str(e)


def _run_pdf_to_pdfa_sync(
    job_id: str,
    in_path: str,
    tmp_dir: str,
    stem: str,
    conformance: str = "PDF/A-2b",
    allow_downgrade: bool = True,
):
    """Convert PDF to PDF/A using Ghostscript or PyMuPDF fallback."""
    try:
        _job_store[job_id]["status"] = JobStatus.PROCESSING
        out_path = os.path.join(tmp_dir, f"{stem}_pdfa.pdf")
        gs_bin = _find_ghostscript()
        converted = False

        # Map conformance level to Ghostscript flags
        conf_upper = (conformance or "").upper().strip()
        if "1" in conf_upper:
            pdfa_flag = "-dPDFA=1"
        elif "3" in conf_upper:
            pdfa_flag = "-dPDFA=3"
        else:
            pdfa_flag = "-dPDFA=2"

        if gs_bin:
            try:
                res = subprocess.run(
                    [
                        gs_bin,
                        pdfa_flag,
                        "-dBATCH",
                        "-dNOPAUSE",
                        "-sColorConversionStrategy=UseDeviceIndependentColor",
                        "-sDEVICE=pdfwrite",
                        f"-sOutputFile={out_path}",
                        in_path,
                    ],
                    capture_output=True,
                    text=True,
                )
                converted = res.returncode == 0 and os.path.isfile(out_path) and os.path.getsize(out_path) > 0
            except Exception:
                converted = False

        if not converted and allow_downgrade:
            try:
                import fitz
                doc = fitz.open(in_path)
                doc.save(out_path, garbage=3, deflate=True)
                doc.close()
                converted = os.path.isfile(out_path) and os.path.getsize(out_path) > 0
            except Exception:
                converted = False

        if not converted:
            _job_store[job_id]["status"] = JobStatus.ERROR
            _job_store[job_id]["error"] = f"PDF/A conversion failed for level {conformance}."
            return

        _job_store[job_id]["status"] = JobStatus.DONE
        _job_store[job_id]["result_path"] = out_path
    except Exception as e:
        _job_store[job_id]["status"] = JobStatus.ERROR
        _job_store[job_id]["error"] = str(e)


def _run_ocr_to_word_sync(job_id: str, in_path: str, tmp_dir: str, stem: str, language: str = "eng"):
    """Perform OCR on an image or PDF file and write extracted text, native tables, logos, and signatures into a Word (.docx) document."""
    try:
        _job_store[job_id]["status"] = JobStatus.PROCESSING
        docx_path = os.path.join(tmp_dir, f"{stem}_ocr.docx")
        doc = Document()

        ext = Path(in_path).suffix.lower()

        img_dir = os.path.join(tmp_dir, "extracted_assets")
        os.makedirs(img_dir, exist_ok=True)

        def _process_single_image(img_path_or_pil, is_scanned_page=False):
            """Process an image file using OpenCV + Tesseract (with graceful PIL fallback) to extract logos, native tables, text, and signatures."""
            try:
                import cv2
                import numpy as np
                CV2_AVAILABLE = True
            except Exception:
                CV2_AVAILABLE = False

            # Load image
            if isinstance(img_path_or_pil, str):
                pil_img = Image.open(img_path_or_pil)
                if CV2_AVAILABLE:
                    cv_img = cv2.imread(img_path_or_pil)
                else:
                    cv_img = None
            else:
                pil_img = img_path_or_pil
                if CV2_AVAILABLE:
                    try:
                        cv_img = cv2.cvtColor(np.array(pil_img.convert('RGB')), cv2.COLOR_RGB2BGR)
                    except Exception:
                        cv_img = None
                else:
                    cv_img = None

            if not CV2_AVAILABLE or cv_img is None:
                if TESSERACT_AVAILABLE:
                    txt = pytesseract.image_to_string(pil_img, lang=language)
                    for block in (txt or "").split('\n\n'):
                        if block.strip():
                            doc.add_paragraph(block.strip())
                return

            img_h, img_w = cv_img.shape[:2]

            # 1. Detect and Crop Top Logo / Header Graphic
            try:
                gray = cv2.cvtColor(cv_img, cv2.COLOR_BGR2GRAY)
                # Look in top 25% of the document for graphics/logos
                top_region = gray[0:int(img_h * 0.28), 0:int(img_w * 0.45)]
                _, top_thresh = cv2.threshold(top_region, 220, 255, cv2.THRESH_BINARY_INV)

                # Find contours in top left region
                top_contours, _ = cv2.findContours(top_thresh, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
                logo_rects = []
                for c in top_contours:
                    x, y, w, h = cv2.boundingRect(c)
                    if w > 35 and h > 35 and (w * h) > 1500:
                        logo_rects.append((x, y, w, h))

                if logo_rects:
                    min_x = min(r[0] for r in logo_rects)
                    min_y = min(r[1] for r in logo_rects)
                    max_x = max(r[0] + r[2] for r in logo_rects)
                    max_y = max(r[1] + r[3] for r in logo_rects)

                    logo_crop = cv_img[max(0, min_y - 5):min(int(img_h * 0.28), max_y + 5),
                                       max(0, min_x - 5):min(int(img_w * 0.45), max_x + 5)]
                    if logo_crop.shape[0] > 20 and logo_crop.shape[1] > 20:
                        logo_file = os.path.join(img_dir, f"logo_{uuid.uuid4().hex[:6]}.png")
                        cv2.imwrite(logo_file, logo_crop)
                        doc.add_picture(logo_file, width=Inches(2.4))
                        doc.add_paragraph() # Spacing
            except Exception as logo_err:
                print(f"Logo detection note: {logo_err}")

            # 2. Detect Table Grid Lines via OpenCV Morphological Operations
            table_cells = []
            try:
                gray = cv2.cvtColor(cv_img, cv2.COLOR_BGR2GRAY)
                _, thresh = cv2.threshold(gray, 200, 255, cv2.THRESH_BINARY_INV)

                # Horizontal line kernel
                h_size = max(20, int(img_w * 0.12))
                h_kernel = cv2.getStructuringElement(cv2.MORPH_RECT, (h_size, 1))
                h_lines = cv2.morphologyEx(thresh, cv2.MORPH_OPEN, h_kernel)

                # Vertical line kernel
                v_size = max(15, int(img_h * 0.03))
                v_kernel = cv2.getStructuringElement(cv2.MORPH_RECT, (1, v_size))
                v_lines = cv2.morphologyEx(thresh, cv2.MORPH_OPEN, v_kernel)

                # Combine horizontal and vertical table grid lines
                table_grid = cv2.add(h_lines, v_lines)
                contours, _ = cv2.findContours(table_grid, cv2.RETR_TREE, cv2.CHAIN_APPROX_SIMPLE)

                for c in contours:
                    x, y, w, h = cv2.boundingRect(c)
                    if w > 40 and h > 14 and w < img_w * 0.95 and h < img_h * 0.8:
                        table_cells.append((x, y, w, h))
            except Exception as tbl_err:
                print(f"Table detection note: {tbl_err}")

            if len(table_cells) >= 4:
                # Group cells into rows and columns
                table_cells.sort(key=lambda c: (c[1], c[0])) # Sort top-to-bottom then left-to-right

                # Row clustering by y coordinate (threshold = 12px)
                rows = []
                for cell in table_cells:
                    x, y, w, h = cell
                    placed = False
                    for r in rows:
                        avg_y = sum(c[1] for c in r) / len(r)
                        if abs(y - avg_y) < 14:
                            r.append(cell)
                            placed = True
                            break
                    if not placed:
                        rows.append([cell])

                # Filter rows with at least 2 cells
                valid_rows = [r for r in rows if len(r) >= 2 or (len(r) == 1 and len(rows) > 3)]
                valid_rows.sort(key=lambda r: sum(c[1] for c in r) / len(r))

                if len(valid_rows) >= 2:
                    max_cols = max(len(r) for r in valid_rows)
                    tbl = doc.add_table(rows=len(valid_rows), cols=max_cols)
                    tbl.style = 'Table Grid'

                    for r_idx, row_cells in enumerate(valid_rows):
                        row_cells.sort(key=lambda c: c[0]) # Left to right
                        for c_idx, (x, y, w, h) in enumerate(row_cells):
                            if c_idx < max_cols:
                                # Crop cell image for OCR
                                pad = 2
                                cell_crop = cv_img[max(0, y+pad):min(img_h, y+h-pad),
                                                   max(0, x+pad):min(img_w, x+w+pad)]
                                if cell_crop.shape[0] > 5 and cell_crop.shape[1] > 5:
                                    cell_pil = Image.fromarray(cv2.cvtColor(cell_crop, cv2.COLOR_BGR2RGB))
                                    cell_text = pytesseract.image_to_string(cell_pil, lang=language, config='--psm 6').strip()
                                    # Clean up OCR noise
                                    cell_text = cell_text.replace('\n', ' ').strip()
                                    tbl.rows[r_idx].cells[c_idx].text = cell_text

                    doc.add_paragraph() # Spacing after table
                else:
                    # Fallback to plain OCR block extraction
                    if TESSERACT_AVAILABLE:
                        txt = pytesseract.image_to_string(pil_img, lang=language)
                        for block in (txt or "").split('\n\n'):
                            if block.strip():
                                doc.add_paragraph(block.strip())
            else:
                # Standard OCR text block extraction
                if TESSERACT_AVAILABLE:
                    txt = pytesseract.image_to_string(pil_img, lang=language)
                    for block in (txt or "").split('\n\n'):
                        if block.strip():
                            doc.add_paragraph(block.strip())

            # 3. Detect and Crop Signature / Stamp Region at Document Bottom
            try:
                bottom_region = cv_img[int(img_h * 0.72):img_h, 0:img_w]
                b_gray = cv2.cvtColor(bottom_region, cv2.COLOR_BGR2GRAY)
                _, b_thresh = cv2.threshold(b_gray, 220, 255, cv2.THRESH_BINARY_INV)

                b_contours, _ = cv2.findContours(b_thresh, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
                sig_rects = []
                for c in b_contours:
                    x, y, w, h = cv2.boundingRect(c)
                    if w > 45 and h > 20 and (w * h) > 1200:
                        sig_rects.append((x, y, w, h))

                if sig_rects:
                    min_x = min(r[0] for r in sig_rects)
                    min_y = min(r[1] for r in sig_rects)
                    max_x = max(r[0] + r[2] for r in sig_rects)
                    max_y = max(r[1] + r[3] for r in sig_rects)

                    sig_crop = bottom_region[max(0, min_y - 10):min(int(img_h * 0.28), max_y + 10),
                                             max(0, min_x - 10):min(img_w, max_x + 10)]
                    if sig_crop.shape[0] > 30 and sig_crop.shape[1] > 50:
                        sig_file = os.path.join(img_dir, f"signature_{uuid.uuid4().hex[:6]}.png")
                        cv2.imwrite(sig_file, sig_crop)
                        doc.add_paragraph() # Spacing
                        doc.add_picture(sig_file, width=Inches(3.8))
            except Exception as sig_err:
                print(f"Signature detection note: {sig_err}")

        if ext == '.pdf':
            import fitz

            fitz_doc = fitz.open(in_path)
            for page_num in range(len(fitz_doc)):
                page = fitz_doc[page_num]
                pix = page.get_pixmap(dpi=150)
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)

                # Process page image with OpenCV logo, table, and signature extraction
                _process_single_image(img, is_scanned_page=True)

                if page_num < len(fitz_doc) - 1:
                    doc.add_page_break()

            fitz_doc.close()

        elif ext in ('.png', '.jpg', '.jpeg', '.webp', '.tiff', '.tif', '.bmp'):
            _process_single_image(in_path, is_scanned_page=False)
        else:
            if PILLOW_AVAILABLE:
                _process_single_image(in_path, is_scanned_page=False)

        doc.save(docx_path)
        _job_store[job_id]["status"] = JobStatus.DONE
        _job_store[job_id]["result_path"] = docx_path
    except Exception as e:
        _job_store[job_id]["status"] = JobStatus.ERROR
        _job_store[job_id]["error"] = str(e)


def _run_pdf_to_markdown_sync(job_id: str, in_path: str, tmp_dir: str, stem: str):
    """Extract structured text and tables from PDF and save as clean Markdown (.md)."""
    try:
        _job_store[job_id]["status"] = JobStatus.PROCESSING
        md_path = os.path.join(tmp_dir, f"{stem}.md")

        import fitz
        doc = fitz.open(in_path)
        md_sections = []

        for page_num in range(len(doc)):
            page = doc[page_num]

            tables = []
            table_rects = []
            if hasattr(page, "find_tables"):
                try:
                    tabs = page.find_tables()
                    for tab in tabs:
                        tables.append(tab)
                        table_rects.append(tab.rect)
                except Exception:
                    pass

            def is_inside_table(bbox):
                rx0, ry0, rx1, ry1 = bbox
                for rect in table_rects:
                    if (rx0 >= rect.x0 - 2 and ry0 >= rect.y0 - 2 and 
                        rx1 <= rect.x1 + 2 and ry1 <= rect.y1 + 2):
                        return True
                return False

            blocks_dict = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)
            blocks = blocks_dict.get("blocks", [])

            font_sizes = []
            for b in blocks:
                if b.get("type") == 0:
                    for line in b.get("lines", []):
                        for span in line.get("spans", []):
                            if span.get("text", "").strip():
                                font_sizes.append(span.get("size", 10.0))

            avg_font_size = sum(font_sizes) / len(font_sizes) if font_sizes else 10.0

            items = []
            for tab in tables:
                items.append(("table", tab.rect.y0, tab))

            for b in blocks:
                if b.get("type") == 0:
                    bbox = b.get("bbox", (0, 0, 0, 0))
                    if not is_inside_table(bbox):
                        items.append(("block", bbox[1], b))

            items.sort(key=lambda x: x[1])

            for item_type, _, obj in items:
                if item_type == "table":
                    try:
                        extracted = obj.extract()
                        if extracted and len(extracted) > 0:
                            header = extracted[0]
                            header_str = "| " + " | ".join(str(c or "").replace("\n", " ").strip() for c in header) + " |"
                            sep_str = "| " + " | ".join("---" for _ in header) + " |"
                            md_sections.append(header_str)
                            md_sections.append(sep_str)
                            for row in extracted[1:]:
                                row_str = "| " + " | ".join(str(c or "").replace("\n", " ").strip() for c in row) + " |"
                                md_sections.append(row_str)
                            md_sections.append("")
                    except Exception:
                        pass
                elif item_type == "block":
                    block_lines = []
                    max_size = 0.0
                    is_bold = False

                    for line in obj.get("lines", []):
                        line_str = ""
                        for span in line.get("spans", []):
                            stext = span.get("text", "")
                            ssize = span.get("size", 10.0)
                            sflags = span.get("flags", 0)
                            sfont = span.get("font", "").lower()
                            if ssize > max_size:
                                max_size = ssize
                            if (sflags & 2) or ("bold" in sfont) or ("black" in sfont) or ("heavy" in sfont):
                                is_bold = True
                            line_str += stext
                        if line_str:
                            block_lines.append(line_str)

                    text_block = "\n".join(block_lines).strip()
                    if not text_block:
                        continue

                    if max_size >= avg_font_size * 1.35:
                        md_sections.append(f"# {text_block}\n")
                    elif max_size >= avg_font_size * 1.18 or (is_bold and max_size > avg_font_size * 1.02 and len(text_block) < 100):
                        md_sections.append(f"## {text_block}\n")
                    elif any(text_block.strip().startswith(prefix) for prefix in ('• ', '- ', '* ', '1. ', '2. ', '3. ')):
                        for line_item in text_block.split('\n'):
                            item_strip = line_item.strip()
                            if item_strip.startswith(('•', '-', '*')):
                                item_content = item_strip.lstrip('•-* ').strip()
                                md_sections.append(f"- {item_content}")
                            else:
                                md_sections.append(item_strip)
                        md_sections.append("")
                    else:
                        md_sections.append(f"{text_block}\n")

        doc.close()

        final_md = "\n".join(md_sections).strip()
        if not final_md:
            final_md = f"# {stem}\n\n(No content extracted)"

        with open(md_path, "w", encoding="utf-8") as f:
            f.write(final_md)

        _job_store[job_id]["status"] = JobStatus.DONE
        _job_store[job_id]["result_path"] = md_path
    except Exception as e:
        _job_store[job_id]["status"] = JobStatus.ERROR
        _job_store[job_id]["error"] = str(e)


def _run_compare_pdf_sync(job_id: str, path1: str, path2: str, tmp_dir: str, stem1: str, stem2: str):
    """Compare two PDF files, extract line/word semantic text diffs, and generate a comparison report PDF."""
    try:
        import fitz
        import difflib
        import html
        import datetime
        from reportlab.lib.pagesizes import letter
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, HRFlowable
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib import colors

        _job_store[job_id]["status"] = JobStatus.PROCESSING

        doc1 = fitz.open(path1)
        doc2 = fitz.open(path2)
        max_pages = max(len(doc1), len(doc2))

        def _extract_page_lines(page):
            lines = []
            if not page:
                return lines
            text_dict = page.get_text("dict")
            for block in text_dict.get("blocks", []):
                if block.get("type") == 0:  # Text block
                    for line in block.get("lines", []):
                        line_text = "".join(span.get("text", "") for span in line.get("spans", [])).strip()
                        if line_text:
                            lines.append({
                                "text": line_text,
                                "bbox": [round(v, 2) for v in line.get("bbox")]
                            })
            return lines

        all_changes = []

        for i in range(max_pages):
            page1 = doc1[i] if i < len(doc1) else None
            page2 = doc2[i] if i < len(doc2) else None

            lines1 = _extract_page_lines(page1)
            lines2 = _extract_page_lines(page2)
            texts1 = [l["text"] for l in lines1]
            texts2 = [l["text"] for l in lines2]

            matcher = difflib.SequenceMatcher(None, texts1, texts2)
            opcodes = matcher.get_opcodes()

            p1_num = i + 1 if page1 else None
            p2_num = i + 1 if page2 else None

            for tag, i1, i2, j1, j2 in opcodes:
                if tag == "equal":
                    continue
                elif tag == "delete":
                    sub_lines = lines1[i1:i2]
                    old_text = "\n".join(l["text"] for l in sub_lines)
                    x0 = min(l["bbox"][0] for l in sub_lines)
                    y0 = min(l["bbox"][1] for l in sub_lines)
                    x1 = max(l["bbox"][2] for l in sub_lines)
                    y1 = max(l["bbox"][3] for l in sub_lines)
                    bbox1 = [round(x0, 2), round(y0, 2), round(x1, 2), round(y1, 2)]
                    all_changes.append({
                        "page1": p1_num,
                        "page2": p2_num,
                        "type": "deletion",
                        "old_text": old_text,
                        "new_text": "",
                        "bbox1": bbox1,
                        "bbox2": None
                    })
                    if page1:
                        for l in sub_lines:
                            # Subtle red underline/outline (no opaque fill blocking text)
                            rect = fitz.Rect(l["bbox"])
                            page1.draw_rect(rect, color=(0.85, 0.2, 0.2), width=1.0, overlay=True)
                            # Strike-through line across text
                            mid_y = (rect.y0 + rect.y1) / 2
                            page1.draw_line(fitz.Point(rect.x0, mid_y), fitz.Point(rect.x1, mid_y), color=(0.85, 0.2, 0.2), width=1.0, overlay=True)
                elif tag == "insert":
                    sub_lines = lines2[j1:j2]
                    new_text = "\n".join(l["text"] for l in sub_lines)
                    x0 = min(l["bbox"][0] for l in sub_lines)
                    y0 = min(l["bbox"][1] for l in sub_lines)
                    x1 = max(l["bbox"][2] for l in sub_lines)
                    y1 = max(l["bbox"][3] for l in sub_lines)
                    bbox2 = [round(x0, 2), round(y0, 2), round(x1, 2), round(y1, 2)]
                    all_changes.append({
                        "page1": p1_num,
                        "page2": p2_num,
                        "type": "addition",
                        "old_text": "",
                        "new_text": new_text,
                        "bbox1": None,
                        "bbox2": bbox2
                    })
                    if page2:
                        for l in sub_lines:
                            # Subtle green outline (no opaque fill blocking text)
                            rect = fitz.Rect(l["bbox"])
                            page2.draw_rect(rect, color=(0.16, 0.65, 0.27), width=1.0, overlay=True)
                elif tag == "replace":
                    old_lines = lines1[i1:i2]
                    new_lines = lines2[j1:j2]
                    old_text = "\n".join(l["text"] for l in old_lines)
                    new_text = "\n".join(l["text"] for l in new_lines)
                    x0_1 = min(l["bbox"][0] for l in old_lines)
                    y0_1 = min(l["bbox"][1] for l in old_lines)
                    x1_1 = max(l["bbox"][2] for l in old_lines)
                    y1_1 = max(l["bbox"][3] for l in old_lines)
                    bbox1 = [round(x0_1, 2), round(y0_1, 2), round(x1_1, 2), round(y1_1, 2)]

                    x0_2 = min(l["bbox"][0] for l in new_lines)
                    y0_2 = min(l["bbox"][1] for l in new_lines)
                    x1_2 = max(l["bbox"][2] for l in new_lines)
                    y1_2 = max(l["bbox"][3] for l in new_lines)
                    bbox2 = [round(x0_2, 2), round(y0_2, 2), round(x1_2, 2), round(y1_2, 2)]

                    all_changes.append({
                        "page1": p1_num,
                        "page2": p2_num,
                        "type": "modification",
                        "old_text": old_text,
                        "new_text": new_text,
                        "bbox1": bbox1,
                        "bbox2": bbox2
                    })
                    if page1:
                        for l in old_lines:
                            # Draw subtle orange side indicator line (gutter bar) or crisp outline instead of solid orange block
                            rect = fitz.Rect(l["bbox"])
                            page1.draw_rect(rect, color=(0.92, 0.45, 0.1), width=1.0, overlay=True)
                    if page2:
                        for l in new_lines:
                            rect = fitz.Rect(l["bbox"])
                            page2.draw_rect(rect, color=(0.92, 0.45, 0.1), width=1.0, overlay=True)

        additions_count = sum(1 for c in all_changes if c["type"] == "addition")
        deletions_count = sum(1 for c in all_changes if c["type"] == "deletion")
        modifications_count = sum(1 for c in all_changes if c["type"] == "modification")
        change_count = len(all_changes)

        comparison_data = {
            "change_count": change_count,
            "additions_count": additions_count,
            "deletions_count": deletions_count,
            "modifications_count": modifications_count,
            "changes": all_changes
        }

        _job_store[job_id]["comparison_data"] = comparison_data

        side_by_side_doc = fitz.open()

        # Standard Landscape Page Dimensions (A4 Landscape: 1190 x 842)
        TARGET_PAGE_W = 1190.0
        TARGET_PAGE_H = 842.0

        for i in range(max_pages):
            pix1 = doc1[i].get_pixmap(dpi=130) if i < len(doc1) else None
            pix2 = doc2[i].get_pixmap(dpi=130) if i < len(doc2) else None

            page = side_by_side_doc.new_page(width=TARGET_PAGE_W, height=TARGET_PAGE_H)

            # Draw Header Bar
            page.draw_rect(fitz.Rect(0, 0, TARGET_PAGE_W, 40), color=None, fill=(0.06, 0.09, 0.16))
            page.insert_text((30, 26), f"DOCUMENT 1 (ORIGINAL): {stem1}.pdf — Page {i+1}", fontsize=11, color=(0.38, 0.65, 0.98))
            page.insert_text((TARGET_PAGE_W / 2 + 20, 26), f"DOCUMENT 2 (MODIFIED): {stem2}.pdf — Page {i+1}", fontsize=11, color=(0.2, 0.8, 0.4))

            # Available canvas area per document side
            side_w = (TARGET_PAGE_W - 60) / 2
            side_h = TARGET_PAGE_H - 60

            if pix1:
                img_path1 = os.path.join(tmp_dir, f"p1_{i}.png")
                pix1.save(img_path1)
                scale = min(side_w / pix1.width, side_h / pix1.height)
                render_w = pix1.width * scale
                render_h = pix1.height * scale
                x_off = 20 + (side_w - render_w) / 2
                y_off = 50 + (side_h - render_h) / 2
                page.insert_image(fitz.Rect(x_off, y_off, x_off + render_w, y_off + render_h), filename=img_path1)

            if pix2:
                img_path2 = os.path.join(tmp_dir, f"p2_{i}.png")
                pix2.save(img_path2)
                scale = min(side_w / pix2.width, side_h / pix2.height)
                render_w = pix2.width * scale
                render_h = pix2.height * scale
                x_off = (TARGET_PAGE_W / 2 + 10) + (side_w - render_w) / 2
                y_off = 50 + (side_h - render_h) / 2
                page.insert_image(fitz.Rect(x_off, y_off, x_off + render_w, y_off + render_h), filename=img_path2)

        doc1.close()
        summary_pdf_path = os.path.join(tmp_dir, "summary_report.pdf")
        doc_rep = SimpleDocTemplate(
            summary_pdf_path,
            pagesize=letter,
            leftMargin=36, rightMargin=36, topMargin=36, bottomMargin=36
        )
        story = []
        styles = getSampleStyleSheet()

        title_style = ParagraphStyle(
            'DocTitle', parent=styles['Heading1'], fontSize=20, leading=24,
            textColor=colors.HexColor('#0F172A'), fontName='Helvetica-Bold', spaceAfter=4
        )
        subtitle_style = ParagraphStyle(
            'DocSubTitle', parent=styles['Normal'], fontSize=10, leading=13,
            textColor=colors.HexColor('#475569'), fontName='Helvetica', spaceAfter=12
        )
        meta_style = ParagraphStyle(
            'MetaText', parent=styles['Normal'], fontSize=9, leading=12,
            textColor=colors.HexColor('#334155'), fontName='Helvetica'
        )
        card_num_style = ParagraphStyle(
            'CardNum', parent=styles['Normal'], fontSize=18, leading=22,
            alignment=1, fontName='Helvetica-Bold'
        )
        card_label_style = ParagraphStyle(
            'CardLabel', parent=styles['Normal'], fontSize=9, leading=11,
            alignment=1, textColor=colors.HexColor('#64748B'), fontName='Helvetica-Bold'
        )
        cell_style = ParagraphStyle(
            'CellText', parent=styles['Normal'], fontSize=8, leading=10,
            textColor=colors.HexColor('#1E293B'), fontName='Helvetica'
        )
        hdr_style = ParagraphStyle(
            'HdrText', parent=styles['Normal'], fontSize=8, leading=10,
            textColor=colors.white, fontName='Helvetica-Bold'
        )

        story.append(Paragraph("PDF Comparison Audit Report", title_style))
        now_str = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        story.append(Paragraph(f"Generated on {now_str} • Side-by-Side Visual & Text Difference Analysis", subtitle_style))
        story.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor('#E2E8F0'), spaceAfter=12))

        meta_data = [
            [
                Paragraph(f"<b>Original (Doc 1):</b> {html.escape(stem1)}.pdf ({max_pages} pages)", meta_style),
                Paragraph(f"<b>Modified (Doc 2):</b> {html.escape(stem2)}.pdf ({max_pages} pages)", meta_style)
            ]
        ]
        meta_table = Table(meta_data, colWidths=[270, 270])
        meta_table.setStyle(TableStyle([
            ('BACKGROUND', (0,0), (-1,-1), colors.HexColor('#F8FAFC')),
            ('PADDING', (0,0), (-1,-1), 8),
            ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
            ('BOX', (0,0), (-1,-1), 0.5, colors.HexColor('#CBD5E1')),
        ]))
        story.append(meta_table)
        story.append(Spacer(1, 12))

        def make_card(num, label, color_hex):
            c_num = ParagraphStyle('CN', parent=card_num_style, textColor=colors.HexColor(color_hex))
            return [Paragraph(str(num), c_num), Paragraph(label, card_label_style)]

        cards_data = [
            [
                make_card(change_count, "TOTAL CHANGES", "#0284C7"),
                make_card(additions_count, "ADDITIONS", "#16A34A"),
                make_card(deletions_count, "DELETIONS", "#DC2626"),
                make_card(modifications_count, "MODIFICATIONS", "#EA580C"),
            ]
        ]
        cards_table = Table(cards_data, colWidths=[135, 135, 135, 135])
        cards_table.setStyle(TableStyle([
            ('BACKGROUND', (0,0), (0,0), colors.HexColor('#F0F9FF')),
            ('BACKGROUND', (1,0), (1,0), colors.HexColor('#F0FDF4')),
            ('BACKGROUND', (2,0), (2,0), colors.HexColor('#FEF2F2')),
            ('BACKGROUND', (3,0), (3,0), colors.HexColor('#FFF7ED')),
            ('BOX', (0,0), (0,0), 0.5, colors.HexColor('#BAE6FD')),
            ('BOX', (1,0), (1,0), 0.5, colors.HexColor('#BBF7D0')),
            ('BOX', (2,0), (2,0), 0.5, colors.HexColor('#FECACA')),
            ('BOX', (3,0), (3,0), 0.5, colors.HexColor('#FFEDD5')),
            ('PADDING', (0,0), (-1,-1), 8),
            ('ALIGN', (0,0), (-1,-1), 'CENTER'),
        ]))
        story.append(cards_table)
        story.append(Spacer(1, 14))

        sec_title = ParagraphStyle('SecTitle', parent=styles['Heading2'], fontSize=12, leading=15, textColor=colors.HexColor('#1E293B'), spaceAfter=8)
        story.append(Paragraph("Detailed Summary Change Log", sec_title))

        if not all_changes:
            no_diff_style = ParagraphStyle('NoDiff', parent=styles['Normal'], fontSize=10, textColor=colors.HexColor('#16A34A'))
            story.append(Paragraph("No differences detected between the two PDF documents.", no_diff_style))
        else:
            table_data = [[
                Paragraph("#", hdr_style),
                Paragraph("Doc 1 Pg", hdr_style),
                Paragraph("Doc 2 Pg", hdr_style),
                Paragraph("Type", hdr_style),
                Paragraph("Original Text (Doc 1)", hdr_style),
                Paragraph("New Text (Doc 2)", hdr_style)
            ]]

            for idx, chg in enumerate(all_changes, 1):
                p1_str = str(chg.get("page1")) if chg.get("page1") is not None else "-"
                p2_str = str(chg.get("page2")) if chg.get("page2") is not None else "-"
                ctype = chg.get("type", "change")

                if ctype == "addition":
                    type_html = "<font color='#16A34A'><b>+ Addition</b></font>"
                elif ctype == "deletion":
                    type_html = "<font color='#DC2626'><b>- Deletion</b></font>"
                else:
                    type_html = "<font color='#EA580C'><b>~ Modified</b></font>"

                old_t = html.escape(chg.get("old_text", "").replace("\n", " "))
                new_t = html.escape(chg.get("new_text", "").replace("\n", " "))
                if len(old_t) > 100:
                    old_t = old_t[:97] + "..."
                if len(new_t) > 100:
                    new_t = new_t[:97] + "..."

                table_data.append([
                    Paragraph(str(idx), cell_style),
                    Paragraph(p1_str, cell_style),
                    Paragraph(p2_str, cell_style),
                    Paragraph(type_html, cell_style),
                    Paragraph(old_t if old_t else "<i>None</i>", cell_style),
                    Paragraph(new_t if new_t else "<i>None</i>", cell_style),
                ])

            change_table = Table(table_data, colWidths=[20, 40, 40, 65, 185, 190])
            change_table.setStyle(TableStyle([
                ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#1E293B')),
                ('PADDING', (0,0), (-1,-1), 4),
                ('VALIGN', (0,0), (-1,-1), 'TOP'),
                ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#E2E8F0')),
                ('ROWBACKGROUNDS', (0,1), (-1,-1), [colors.white, colors.HexColor('#F8FAFC')]),
            ]))
            story.append(change_table)

        doc_rep.build(story)

        summary_doc = fitz.open(summary_pdf_path)
        final_doc = fitz.open()

        final_doc.insert_pdf(summary_doc)
        summary_doc.close()

        final_doc.insert_pdf(side_by_side_doc)
        side_by_side_doc.close()

        out_path = os.path.join(tmp_dir, f"{stem1}_vs_{stem2}_comparison.pdf")
        final_doc.save(out_path, garbage=3, deflate=True)
        final_doc.close()

        _job_store[job_id]["status"] = JobStatus.DONE
        _job_store[job_id]["result_path"] = out_path
    except Exception as e:
        _job_store[job_id]["status"] = JobStatus.ERROR
        _job_store[job_id]["error"] = str(e)


# ─── Helper tools ─────────────────────────────────────────────────────────

def _find_libreoffice():
    """Find the LibreOffice soffice executable."""
    candidates = [
        r'C:\Program Files\LibreOffice\program\soffice.exe',
        r'C:\Program Files (x86)\LibreOffice\program\soffice.exe',
        shutil.which('soffice'),
    ]
    for c in candidates:
        if c and os.path.isfile(c):
            return c
    return None


def _ocr_pdf_to_text(pdf_bytes: bytes) -> str:
    """Use pytesseract to OCR a scanned PDF."""
    if not TESSERACT_AVAILABLE or not PDF2IMAGE_AVAILABLE:
        return ""
    try:
        images = convert_from_bytes(pdf_bytes, dpi=300, poppler_path=POPPLER_PATH)
        text_parts = []
        for img in images:
            page_text = pytesseract.image_to_string(img)
            text_parts.append(page_text)
        return "\n\n".join(text_parts)
    except Exception as e:
        print(f"   [!] OCR failed: {e}")
        return ""


def _ocr_pdf_to_docx(pdf_bytes: bytes, docx_path: str):
    """
    OCR a scanned PDF and create a DOCX that preserves table structure.
    Uses PyMuPDF's built-in OCR.
    """
    import fitz  # PyMuPDF
    from docx.shared import Pt

    env_path = os.environ.get('PATH', '')
    tess_dir = r'C:\Program Files\Tesseract-OCR'
    if os.path.isdir(tess_dir) and tess_dir not in env_path:
        os.environ['PATH'] = tess_dir + ';' + env_path

    try:
        pdf_doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        print(f"   [i] Opened PDF: {len(pdf_doc)} pages")

        doc = Document()
        style = doc.styles['Normal']
        style.font.size = Pt(10)
        style.font.name = 'Calibri'

        for page_idx in range(len(pdf_doc)):
            page = pdf_doc[page_idx]
            if page_idx > 0:
                doc.add_page_break()

            print(f"   [i] OCR'ing page {page_idx + 1}/{len(pdf_doc)}...")

            tp = page.get_textpage_ocr(language="eng", dpi=300, full=True)
            words = page.get_text("words", textpage=tp)
            print(f"   [i]   OCR found {len(words)} words")

            if not words:
                continue

            page_width = page.rect.width

            # Group words into visual lines
            line_tolerance = 5
            lines = []
            current_line = [words[0]]

            for w in words[1:]:
                avg_y = sum(ww[1] for ww in current_line) / len(current_line)
                if abs(w[1] - avg_y) <= line_tolerance:
                    current_line.append(w)
                else:
                    current_line.sort(key=lambda x: x[0])
                    lines.append(current_line)
                    current_line = [w]
            if current_line:
                current_line.sort(key=lambda x: x[0])
                lines.append(current_line)

            # Detect column structure
            min_gap = page_width * 0.03
            all_gap_midpoints = []

            for line in lines:
                if len(line) < 2:
                    continue
                for i in range(len(line) - 1):
                    gap = line[i + 1][0] - line[i][2]
                    if gap > min_gap:
                        midpoint = (line[i][2] + line[i + 1][0]) / 2
                        all_gap_midpoints.append(midpoint)

            strong_boundaries = []
            if all_gap_midpoints:
                all_gap_midpoints.sort()
                cluster_tolerance = page_width * 0.04
                clusters = []
                current_cluster = [all_gap_midpoints[0]]

                for gm in all_gap_midpoints[1:]:
                    if gm - current_cluster[-1] <= cluster_tolerance:
                        current_cluster.append(gm)
                    else:
                        clusters.append(current_cluster)
                        current_cluster = [gm]
                clusters.append(current_cluster)

                min_appearances = min(3, max(2, len(lines) // 5))
                for cluster in clusters:
                    if len(cluster) >= min_appearances:
                        strong_boundaries.append(sorted(cluster)[len(cluster) // 2])

            if strong_boundaries:
                num_cols = len(strong_boundaries) + 1
                col_ranges = []
                prev = 0
                for b in strong_boundaries:
                    col_ranges.append((prev, b))
                    prev = b
                col_ranges.append((prev, page_width + 50))

                def assign_to_columns(line_words):
                    row = [''] * num_cols
                    for w in line_words:
                        word_center = (w[0] + w[2]) / 2
                        for ci, (cl, cr) in enumerate(col_ranges):
                            if cl <= word_center < cr:
                                row[ci] = (row[ci] + ' ' + w[4]).strip() if row[ci] else w[4]
                                break
                        else:
                            row[-1] = (row[-1] + ' ' + w[4]).strip() if row[-1] else w[4]
                    return row

                min_cols_for_table = min(3, num_cols)
                blocks = []
                pending_text = []
                pending_table = []

                for line in lines:
                    row = assign_to_columns(line)
                    populated = sum(1 for c in row if c.strip())

                    if populated >= min_cols_for_table:
                        if pending_text:
                            blocks.append(('text', pending_text))
                            pending_text = []
                        pending_table.append(row)
                    else:
                        if pending_table:
                            blocks.append(('table', pending_table))
                            pending_table = []
                        line_text = ' '.join(w[4] for w in line)
                        if line_text.strip():
                            pending_text.append(line_text.strip())

                if pending_table:
                    blocks.append(('table', pending_table))
                if pending_text:
                    blocks.append(('text', pending_text))

                for btype, bdata in blocks:
                    if btype == 'text':
                        for text_line in bdata:
                            if text_line:
                                p = doc.add_paragraph(text_line)
                                if text_line.isupper() and len(text_line) > 3:
                                    for run in p.runs:
                                        run.bold = True
                    elif btype == 'table':
                        if not bdata:
                            continue
                        max_col = 0
                        for row in bdata:
                            for ci in range(len(row) - 1, -1, -1):
                                if row[ci].strip():
                                    max_col = max(max_col, ci + 1)
                                    break
                        if max_col < 2:
                            for row in bdata:
                                text = ' '.join(c for c in row if c.strip())
                                if text:
                                    doc.add_paragraph(text)
                            continue

                        table = doc.add_table(rows=len(bdata), cols=max_col)
                        table.style = 'Table Grid'
                        for r_idx, row_data in enumerate(bdata):
                            for c_idx in range(max_col):
                                cell = table.cell(r_idx, c_idx)
                                cell.text = (row_data[c_idx] if c_idx < len(row_data) else '').strip()
                                if r_idx == 0:
                                    for p in cell.paragraphs:
                                        for run in p.runs:
                                            run.bold = True
                        doc.add_paragraph('')
            else:
                for line in lines:
                    line_text = ' '.join(w[4] for w in line)
                    if line_text.strip():
                        p = doc.add_paragraph(line_text.strip())
                        if line_text.strip().isupper() and len(line_text.strip()) > 3:
                            for run in p.runs:
                                run.bold = True

        pdf_doc.close()
        doc.save(docx_path)

        final_doc = Document(docx_path)
        total_text = "\n".join(p.text for p in final_doc.paragraphs).strip()
        total_tables = len(final_doc.tables)
        print(f"   [✓] OCR complete: {len(total_text)} chars, {total_tables} tables")

        return True

    except Exception as e:
        print(f"   [!] PyMuPDF OCR failed: {e}")
        import traceback
        traceback.print_exc()

        if TESSERACT_AVAILABLE and PDF2IMAGE_AVAILABLE:
            try:
                images = convert_from_bytes(pdf_bytes, dpi=300, poppler_path=POPPLER_PATH)
                doc = Document()
                for page_idx, img in enumerate(images):
                    if page_idx > 0:
                        doc.add_page_break()
                    page_text = pytesseract.image_to_string(img)
                    for para in page_text.split('\n'):
                        if para.strip():
                            doc.add_paragraph(para.strip())
                doc.save(docx_path)
                return True
            except:
                pass
        return False


# ─── Endpoints ────────────────────────────────────────────────────────────

@router.get("/api/health")
async def health_check():
    return {"status": "ok"}


@router.get("/api/convert/capabilities")
async def conversion_capabilities():
    """Return which conversion tools are available on this server."""
    libre = _find_libreoffice()
    return {
        "pdf_to_word": PDF2DOCX_AVAILABLE,
        "word_to_pdf": libre is not None,
        "pdf_to_text": True,
        "image_to_pdf": PILLOW_AVAILABLE,
        "pdf_to_images": PDF2IMAGE_AVAILABLE,
        "merge_pdf": True,
        "compress_pdf": True,
        "ocr": TESSERACT_AVAILABLE and PDF2IMAGE_AVAILABLE,
        "libreoffice_path": libre,
    }





@router.post("/api/convert/pdf-to-word")
async def convert_pdf_to_word(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
):
    """Submit a PDF→DOCX job. Returns job_id immediately; conversion runs in background."""
    if not PDF2DOCX_AVAILABLE:
        raise HTTPException(status_code=503, detail="pdf2docx is not installed on this server.")
    if not file.filename.lower().endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Please upload a PDF file.")

    stem = Path(file.filename).stem
    tmp_dir = tempfile.mkdtemp()
    in_path = os.path.join(tmp_dir, "input.pdf")
    CHUNK = 1024 * 1024
    try:
        with open(in_path, "wb") as out_f:
            while True:
                chunk = await file.read(CHUNK)
                if not chunk:
                    break
                out_f.write(chunk)
    except Exception as e:
        shutil.rmtree(tmp_dir, ignore_errors=True)
        raise HTTPException(status_code=500, detail=f"Upload failed: {e}")

    with open(in_path, "rb") as f:
        if f.read(4) != b"%PDF":
            shutil.rmtree(tmp_dir, ignore_errors=True)
            raise HTTPException(status_code=400, detail="Not a valid PDF file.")

    job_id = str(uuid.uuid4())
    _job_store[job_id] = {
        "status": JobStatus.PENDING,
        "result_path": None,
        "tmp_dir": tmp_dir,
        "filename": f"{stem}.docx",
        "content_type": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "error": None,
    }
    background_tasks.add_task(
        asyncio.get_event_loop().run_in_executor,
        _executor,
        _run_pdf_to_word_sync,
        job_id, in_path, tmp_dir,
    )
    asyncio.create_task(_schedule_job_cleanup(job_id))
    return JSONResponse(status_code=202, content={"job_id": job_id, "status": JobStatus.PENDING})


@router.post("/api/convert/word-to-pdf")
async def convert_word_to_pdf(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
):
    """Submit a Word→PDF job. Returns job_id immediately; conversion runs in background."""
    soffice = _find_libreoffice()
    if not soffice:
        raise HTTPException(status_code=503, detail="LibreOffice is not installed.")

    ext = Path(file.filename).suffix.lower()
    if ext not in (".docx", ".doc", ".odt", ".rtf"):
        raise HTTPException(
            status_code=400,
            detail="Please upload a Word document (.docx, .doc, .odt, or .rtf).",
        )

    stem = Path(file.filename).stem
    tmp_dir = tempfile.mkdtemp()
    in_path = os.path.join(tmp_dir, f"{stem}{ext}")  # preserve original name for LibreOffice output
    CHUNK = 1024 * 1024
    try:
        with open(in_path, "wb") as out_f:
            while True:
                chunk = await file.read(CHUNK)
                if not chunk:
                    break
                out_f.write(chunk)
    except Exception as e:
        shutil.rmtree(tmp_dir, ignore_errors=True)
        raise HTTPException(status_code=500, detail=f"Upload failed: {e}")

    job_id = str(uuid.uuid4())
    _job_store[job_id] = {
        "status": JobStatus.PENDING,
        "result_path": None,
        "tmp_dir": tmp_dir,
        "filename": f"{stem}.pdf",
        "content_type": "application/pdf",
        "error": None,
    }
    background_tasks.add_task(
        asyncio.get_event_loop().run_in_executor,
        _executor,
        _run_word_to_pdf_sync,
        job_id, in_path, tmp_dir, soffice,
    )
    asyncio.create_task(_schedule_job_cleanup(job_id))
    return JSONResponse(status_code=202, content={"job_id": job_id, "status": JobStatus.PENDING})


@router.get("/api/jobs/{job_id}/status")
async def job_status(job_id: str):
    """Universal status poll for any async conversion job."""
    job = _job_store.get(job_id)
    if not job:
        raise HTTPException(status_code=404, detail="Job not found.")
    status_str = job["status"].value if hasattr(job["status"], "value") else str(job["status"])
    response = {"job_id": job_id, "status": status_str}
    if job["status"] == JobStatus.ERROR:
        response["error"] = job.get("error", "Unknown error")
    if job["status"] == JobStatus.DONE:
        response["filename"] = job.get("filename", "output")
        if "comparison_data" in job and job["comparison_data"] is not None:
            response["comparison_data"] = job["comparison_data"]
    return JSONResponse(content=response)


@router.post("/api/jobs/{job_id}/cancel")
async def cancel_job(job_id: str):
    """Cancel a running or pending conversion job on the backend."""
    job = _job_store.get(job_id)
    if not job:
        return JSONResponse(content={"status": "not_found", "message": "Job not found or already cleaned up."})
    
    job["status"] = JobStatus.ERROR
    job["error"] = "Cancelled by user."
    job["cancelled"] = True

    proc = job.get("process")
    if proc:
        try:
            proc.terminate()
            proc.kill()
        except Exception:
            pass

    tmp_dir = job.get("tmp_dir")
    if tmp_dir and os.path.exists(tmp_dir):
        try:
            shutil.rmtree(tmp_dir, ignore_errors=True)
        except Exception:
            pass

    return JSONResponse(content={"status": "cancelled", "job_id": job_id})



@router.get("/api/jobs/{job_id}/download")
async def job_download(job_id: str):
    """Universal download for any async conversion job. File kept for 30 min (re-downloadable)."""
    job = _job_store.get(job_id)
    if not job:
        raise HTTPException(status_code=404, detail="Job not found or expired (30-min window).")
    if job["status"] != JobStatus.DONE:
        raise HTTPException(status_code=409, detail="Conversion not complete yet.")

    result_path = job.get("result_path")
    if not result_path or not os.path.isfile(result_path):
        raise HTTPException(status_code=500, detail="Output file not found.")

    extra_headers = job.get("extra_headers", {})
    return FileResponse(
        path=result_path,
        media_type=job.get("content_type", "application/octet-stream"),
        filename=job.get("filename", Path(result_path).name),
        headers=extra_headers,
    )


@router.post("/api/convert/pdf-to-text")
async def convert_pdf_to_text(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
):
    """Submit a PDF→TXT job. Returns job_id immediately; extraction runs in background."""
    if not file.filename.lower().endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Please upload a PDF file.")

    stem = Path(file.filename).stem
    tmp_dir = tempfile.mkdtemp()
    in_path = os.path.join(tmp_dir, "input.pdf")
    CHUNK = 1024 * 1024
    try:
        with open(in_path, "wb") as out_f:
            while True:
                chunk = await file.read(CHUNK)
                if not chunk:
                    break
                out_f.write(chunk)
    except Exception as e:
        shutil.rmtree(tmp_dir, ignore_errors=True)
        raise HTTPException(status_code=500, detail=f"Upload failed: {e}")

    with open(in_path, "rb") as f:
        if f.read(4) != b"%PDF":
            shutil.rmtree(tmp_dir, ignore_errors=True)
            raise HTTPException(status_code=400, detail="Not a valid PDF file.")

    job_id = str(uuid.uuid4())
    _job_store[job_id] = {
        "status": JobStatus.PENDING,
        "result_path": None,
        "tmp_dir": tmp_dir,
        "filename": f"{stem}.txt",
        "content_type": "text/plain; charset=utf-8",
        "error": None,
    }
    background_tasks.add_task(
        asyncio.get_event_loop().run_in_executor,
        _executor,
        _run_pdf_to_text_sync,
        job_id, in_path, tmp_dir, stem,
    )
    asyncio.create_task(_schedule_job_cleanup(job_id))
    return JSONResponse(status_code=202, content={"job_id": job_id, "status": JobStatus.PENDING})


@router.post("/api/convert/image-to-pdf")
async def convert_image_to_pdf(
    background_tasks: BackgroundTasks,
    files: List[UploadFile] = File(...),
):
    """Submit an image→PDF job. Returns job_id immediately."""
    if not PILLOW_AVAILABLE:
        raise HTTPException(status_code=503, detail="Pillow is not installed.")

    tmp_dir = tempfile.mkdtemp()
    file_paths = []
    CHUNK = 1024 * 1024
    for f in files:
        ext = Path(f.filename).suffix.lower()
        if ext not in ('.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.tif', '.webp'):
            shutil.rmtree(tmp_dir, ignore_errors=True)
            raise HTTPException(status_code=400, detail=f"Unsupported image format: {f.filename}.")
        in_path = os.path.join(tmp_dir, Path(f.filename).name)
        try:
            with open(in_path, "wb") as out_f:
                while True:
                    chunk = await f.read(CHUNK)
                    if not chunk:
                        break
                    out_f.write(chunk)
        except Exception as e:
            shutil.rmtree(tmp_dir, ignore_errors=True)
            raise HTTPException(status_code=500, detail=f"Upload failed: {e}")
        file_paths.append(in_path)

    if not file_paths:
        shutil.rmtree(tmp_dir, ignore_errors=True)
        raise HTTPException(status_code=400, detail="No valid images provided.")

    first_stem = Path(files[0].filename).stem
    out_filename = f"{first_stem}.pdf" if len(files) == 1 else f"{first_stem}_and_{len(files) - 1}_more.pdf"
    job_id = str(uuid.uuid4())
    _job_store[job_id] = {
        "status": JobStatus.PENDING, "result_path": None, "tmp_dir": tmp_dir,
        "filename": out_filename, "content_type": "application/pdf", "error": None,
    }
    background_tasks.add_task(
        asyncio.get_event_loop().run_in_executor, _executor,
        _run_image_to_pdf_sync, job_id, file_paths, tmp_dir,
    )
    asyncio.create_task(_schedule_job_cleanup(job_id))
    return JSONResponse(status_code=202, content={"job_id": job_id, "status": JobStatus.PENDING})


@router.post("/api/convert/pdf-to-images")
async def convert_pdf_to_images(
    background_tasks: BackgroundTasks,
    files: List[UploadFile] = File(...),
):
    """Submit one or more PDFs→images ZIP job. Accepts multiple PDFs simultaneously.
    All images are packed into a single ZIP (sub-folders per PDF when >1 file).
    """
    if not PDF2IMAGE_AVAILABLE:
        raise HTTPException(status_code=503, detail="pdf2image is not installed.")
    if not files:
        raise HTTPException(status_code=400, detail="Please upload at least one PDF file.")

    CHUNK = 1024 * 1024
    tmp_dir = tempfile.mkdtemp()
    in_paths = []
    stems = []

    for upload in files:
        if not upload.filename.lower().endswith(".pdf"):
            shutil.rmtree(tmp_dir, ignore_errors=True)
            raise HTTPException(status_code=400, detail=f"'{upload.filename}' is not a PDF file.")
        stem = Path(upload.filename).stem
        in_path = os.path.join(tmp_dir, f"{stem}_{len(in_paths)}.pdf")
        try:
            with open(in_path, "wb") as out_f:
                while True:
                    chunk = await upload.read(CHUNK)
                    if not chunk:
                        break
                    out_f.write(chunk)
        except Exception as e:
            shutil.rmtree(tmp_dir, ignore_errors=True)
            raise HTTPException(status_code=500, detail=f"Upload failed for '{upload.filename}': {e}")
        with open(in_path, "rb") as vf:
            if vf.read(4) != b"%PDF":
                shutil.rmtree(tmp_dir, ignore_errors=True)
                raise HTTPException(status_code=400, detail=f"'{upload.filename}' is not a valid PDF.")
        in_paths.append(in_path)
        stems.append(stem)

    zip_name = f"{stems[0]}_images.zip" if len(stems) == 1 else "pdf_images.zip"
    job_id = str(uuid.uuid4())
    _job_store[job_id] = {
        "status": JobStatus.PENDING,
        "result_path": None,
        "tmp_dir": tmp_dir,
        "filename": zip_name,
        "content_type": "application/zip",
        "error": None,
    }
    _executor.submit(
        _run_pdf_to_images_sync,
        job_id, in_paths, tmp_dir, stems,
    )
    asyncio.create_task(_schedule_job_cleanup(job_id))
    return JSONResponse(status_code=202, content={"job_id": job_id, "status": JobStatus.PENDING})


@router.post("/api/convert/merge-pdf")
async def merge_pdfs(
    background_tasks: BackgroundTasks,
    files: List[UploadFile] = File(...),
):
    """Submit a merge-PDF job. Returns job_id immediately."""
    if len(files) < 2:
        raise HTTPException(status_code=400, detail="Please upload at least 2 PDF files to merge.")

    tmp_dir = tempfile.mkdtemp()
    file_paths = []
    CHUNK = 1024 * 1024
    for i, f in enumerate(files):
        if not f.filename.lower().endswith('.pdf'):
            shutil.rmtree(tmp_dir, ignore_errors=True)
            raise HTTPException(status_code=400, detail=f"All files must be PDFs. '{f.filename}' is not a PDF.")
        in_path = os.path.join(tmp_dir, f"{i}_{Path(f.filename).name}")
        try:
            with open(in_path, "wb") as out_f:
                while True:
                    chunk = await f.read(CHUNK)
                    if not chunk:
                        break
                    out_f.write(chunk)
        except Exception as e:
            shutil.rmtree(tmp_dir, ignore_errors=True)
            raise HTTPException(status_code=500, detail=f"Upload failed: {e}")
        with open(in_path, "rb") as vf:
            if b'%PDF' not in vf.read(1024):
                shutil.rmtree(tmp_dir, ignore_errors=True)
                raise HTTPException(status_code=400, detail=f"'{f.filename}' is not a valid PDF.")
        file_paths.append(in_path)

    first_stem = Path(files[0].filename).stem
    job_id = str(uuid.uuid4())
    _job_store[job_id] = {
        "status": JobStatus.PENDING,
        "result_path": None,
        "tmp_dir": tmp_dir,
        "filename": f"{first_stem}_merged.pdf",
        "content_type": "application/pdf",
        "error": None,
    }
    background_tasks.add_task(
        asyncio.get_event_loop().run_in_executor,
        _executor,
        _run_merge_pdfs_sync,
        job_id, file_paths, tmp_dir,
    )
    asyncio.create_task(_schedule_job_cleanup(job_id))
    return JSONResponse(status_code=202, content={"job_id": job_id, "status": JobStatus.PENDING})


@router.post("/api/convert/compress-pdf")
async def compress_pdf(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    compression_level: str = Form("recommended"),
    dpi: int = Form(150),
    quality: int = Form(65),
):
    """Submit a compress-PDF job. Returns job_id immediately."""
    if not file.filename.lower().endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Please upload a PDF file.")

    stem = Path(file.filename).stem
    tmp_dir = tempfile.mkdtemp()
    in_path = os.path.join(tmp_dir, "input.pdf")
    CHUNK = 1024 * 1024
    try:
        with open(in_path, "wb") as out_f:
            while True:
                chunk = await file.read(CHUNK)
                if not chunk:
                    break
                out_f.write(chunk)
    except Exception as e:
        shutil.rmtree(tmp_dir, ignore_errors=True)
        raise HTTPException(status_code=500, detail=f"Upload failed: {e}")

    with open(in_path, "rb") as f:
        if f.read(4) != b'%PDF':
            shutil.rmtree(tmp_dir, ignore_errors=True)
            raise HTTPException(status_code=400, detail="Not a valid PDF file.")

    job_id = str(uuid.uuid4())
    _job_store[job_id] = {
        "status": JobStatus.PENDING,
        "result_path": None,
        "tmp_dir": tmp_dir,
        "filename": file.filename,
        "content_type": "application/pdf",
        "error": None,
    }
    background_tasks.add_task(
        asyncio.get_event_loop().run_in_executor,
        _executor,
        _run_compress_pdf_sync,
        job_id, in_path, tmp_dir, compression_level, int(dpi), int(quality),
    )
    asyncio.create_task(_schedule_job_cleanup(job_id))
    return JSONResponse(status_code=202, content={"job_id": job_id, "status": JobStatus.PENDING})


@router.post("/api/convert/compress-image")
async def compress_image(
    background_tasks: BackgroundTasks,
    files: List[UploadFile] = File(...),
    quality: int = Form(75),
    target_format: str = Form("original"),
    max_dim: int = Form(0),
):
    """Submit an image compression job. Returns job_id immediately."""
    if not PILLOW_AVAILABLE:
        raise HTTPException(status_code=503, detail="Pillow is not installed.")

    if not files:
        raise HTTPException(status_code=400, detail="Please upload at least one image file.")

    ALLOWED_EXTS = {".jpg", ".jpeg", ".png", ".webp", ".bmp", ".tiff", ".tif"}
    tmp_dir = tempfile.mkdtemp()
    file_paths = []
    CHUNK = 1024 * 1024

    for idx, f in enumerate(files):
        ext = Path(f.filename or "").suffix.lower()
        if ext not in ALLOWED_EXTS:
            shutil.rmtree(tmp_dir, ignore_errors=True)
            raise HTTPException(
                status_code=400,
                detail=f"Unsupported image format for '{f.filename}'. Supported formats: .jpg, .jpeg, .png, .webp, .bmp, .tiff, .tif.",
            )
        orig_name = Path(f.filename or f"image_{idx}{ext}").name
        in_path = os.path.join(tmp_dir, f"{idx}___{orig_name}")
        try:
            with open(in_path, "wb") as out_f:
                while True:
                    chunk = await f.read(CHUNK)
                    if not chunk:
                        break
                    out_f.write(chunk)
        except Exception as e:
            shutil.rmtree(tmp_dir, ignore_errors=True)
            raise HTTPException(status_code=500, detail=f"Upload failed: {e}")
        file_paths.append(in_path)

    job_id = str(uuid.uuid4())
    _job_store[job_id] = {
        "status": JobStatus.PENDING,
        "result_path": None,
        "tmp_dir": tmp_dir,
        "filename": "compressed_output",
        "content_type": "image/jpeg",
        "error": None,
    }
    background_tasks.add_task(
        asyncio.get_event_loop().run_in_executor,
        _executor,
        _run_compress_image_sync,
        job_id, file_paths, tmp_dir, quality, target_format, max_dim,
    )
    asyncio.create_task(_schedule_job_cleanup(job_id))
    return JSONResponse(status_code=202, content={"job_id": job_id, "status": JobStatus.PENDING})


@router.post("/api/convert/split-pdf")
async def split_pdf(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    ranges: str = Form("1"),
    merge_ranges: bool = Form(False),
    max_size_kb: int = Form(0),
):
    """Submit a split-PDF job. ranges = semicolon-separated page groups, e.g. '1-3 ; 4-6 ; 7'"""
    if not file.filename.lower().endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Please upload a PDF file.")

    stem = Path(file.filename).stem
    tmp_dir = tempfile.mkdtemp()
    in_path = os.path.join(tmp_dir, "input.pdf")
    CHUNK = 1024 * 1024
    try:
        with open(in_path, "wb") as out_f:
            while True:
                chunk = await file.read(CHUNK)
                if not chunk:
                    break
                out_f.write(chunk)
    except Exception as e:
        shutil.rmtree(tmp_dir, ignore_errors=True)
        raise HTTPException(status_code=500, detail=f"Upload failed: {e}")
    with open(in_path, "rb") as f:
        if f.read(4) != b'%PDF':
            shutil.rmtree(tmp_dir, ignore_errors=True)
            raise HTTPException(status_code=400, detail="Not a valid PDF file.")

    job_id = str(uuid.uuid4())
    _job_store[job_id] = {
        "status": JobStatus.PENDING, "result_path": None, "tmp_dir": tmp_dir,
        "filename": f"{stem}_split.zip", "content_type": "application/zip", "error": None,
    }
    background_tasks.add_task(
        asyncio.get_event_loop().run_in_executor, _executor,
        _run_split_pdf_sync, job_id, in_path, tmp_dir, stem, ranges, merge_ranges, max_size_kb,
    )
    asyncio.create_task(_schedule_job_cleanup(job_id))
    return JSONResponse(status_code=202, content={"job_id": job_id, "status": JobStatus.PENDING})


@router.post("/api/convert/remove-pages")
async def remove_pages(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    pages: str = Form("1"),
):
    """Remove specified pages from a PDF. pages = comma-separated page numbers, e.g. '2, 5, 8'"""
    if not file.filename.lower().endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Please upload a PDF file.")

    stem = Path(file.filename).stem
    tmp_dir = tempfile.mkdtemp()
    in_path = os.path.join(tmp_dir, "input.pdf")
    CHUNK = 1024 * 1024
    try:
        with open(in_path, "wb") as out_f:
            while True:
                chunk = await file.read(CHUNK)
                if not chunk:
                    break
                out_f.write(chunk)
    except Exception as e:
        shutil.rmtree(tmp_dir, ignore_errors=True)
        raise HTTPException(status_code=500, detail=f"Upload failed: {e}")
    with open(in_path, "rb") as f:
        if f.read(4) != b'%PDF':
            shutil.rmtree(tmp_dir, ignore_errors=True)
            raise HTTPException(status_code=400, detail="Not a valid PDF file.")

    job_id = str(uuid.uuid4())
    _job_store[job_id] = {
        "status": JobStatus.PENDING, "result_path": None, "tmp_dir": tmp_dir,
        "filename": f"{stem}_removed.pdf", "content_type": "application/pdf", "error": None,
    }
    background_tasks.add_task(
        asyncio.get_event_loop().run_in_executor, _executor,
        _run_remove_pages_sync, job_id, in_path, tmp_dir, stem, pages,
    )
    asyncio.create_task(_schedule_job_cleanup(job_id))
    return JSONResponse(status_code=202, content={"job_id": job_id, "status": JobStatus.PENDING})


@router.post("/api/convert/extract-pages")
async def extract_pages(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    pages: str = Form("1"),
):
    """Extract specific pages from a PDF. pages = comma/range string, e.g. '1-3, 5'"""
    if not file.filename.lower().endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Please upload a PDF file.")

    stem = Path(file.filename).stem
    tmp_dir = tempfile.mkdtemp()
    in_path = os.path.join(tmp_dir, "input.pdf")
    CHUNK = 1024 * 1024
    try:
        with open(in_path, "wb") as out_f:
            while True:
                chunk = await file.read(CHUNK)
                if not chunk:
                    break
                out_f.write(chunk)
    except Exception as e:
        shutil.rmtree(tmp_dir, ignore_errors=True)
        raise HTTPException(status_code=500, detail=f"Upload failed: {e}")
    with open(in_path, "rb") as f:
        if f.read(4) != b'%PDF':
            shutil.rmtree(tmp_dir, ignore_errors=True)
            raise HTTPException(status_code=400, detail="Not a valid PDF file.")

    job_id = str(uuid.uuid4())
    _job_store[job_id] = {
        "status": JobStatus.PENDING, "result_path": None, "tmp_dir": tmp_dir,
        "filename": f"{stem}_extracted.pdf", "content_type": "application/pdf", "error": None,
    }
    background_tasks.add_task(
        asyncio.get_event_loop().run_in_executor, _executor,
        _run_extract_pages_sync, job_id, in_path, tmp_dir, stem, pages,
    )
    asyncio.create_task(_schedule_job_cleanup(job_id))
    return JSONResponse(status_code=202, content={"job_id": job_id, "status": JobStatus.PENDING})


@router.post("/api/convert/organize-pdf")
async def organize_pdf(
    background_tasks: BackgroundTasks,
    files: Optional[List[UploadFile]] = File(None),
    file: Optional[UploadFile] = File(None),
    order: str = Form("0:0, 0:1, 1:0"),
    add_page_numbers: bool = Form(False),
    page_number_position: str = Form("bottom-center"),
    page_number_format: str = Form("Page {page} of {total}"),
    start_number: int = Form(1),
):
    """Reorder and assemble PDF pages from single or multiple PDFs. order = sequence of file_index:page_index or page indices."""
    upload_files = []
    if files:
        upload_files.extend(files)
    if file and file not in upload_files:
        upload_files.append(file)

    if not upload_files:
        raise HTTPException(status_code=400, detail="Please upload at least one PDF file.")

    tmp_dir = tempfile.mkdtemp()
    file_paths = []
    CHUNK = 1024 * 1024
    stem = "organized"

    for idx, f in enumerate(upload_files):
        if not f.filename or not f.filename.lower().endswith('.pdf'):
            shutil.rmtree(tmp_dir, ignore_errors=True)
            raise HTTPException(status_code=400, detail=f"File '{f.filename}' is not a valid PDF file.")

        if idx == 0:
            stem = Path(f.filename).stem

        in_path = os.path.join(tmp_dir, f"input_{idx}.pdf")
        try:
            with open(in_path, "wb") as out_f:
                while True:
                    chunk = await f.read(CHUNK)
                    if not chunk:
                        break
                    out_f.write(chunk)
        except Exception as e:
            shutil.rmtree(tmp_dir, ignore_errors=True)
            raise HTTPException(status_code=500, detail=f"Upload failed: {e}")

        with open(in_path, "rb") as check_f:
            if check_f.read(4) != b'%PDF':
                shutil.rmtree(tmp_dir, ignore_errors=True)
                raise HTTPException(status_code=400, detail=f"File '{f.filename}' is not a valid PDF file.")

        file_paths.append(in_path)

    job_id = str(uuid.uuid4())
    _job_store[job_id] = {
        "status": JobStatus.PENDING, "result_path": None, "tmp_dir": tmp_dir,
        "filename": f"{stem}_organized.pdf", "content_type": "application/pdf", "error": None,
    }
    background_tasks.add_task(
        asyncio.get_event_loop().run_in_executor, _executor,
        _run_organize_pdf_sync, job_id, file_paths, tmp_dir, stem, order,
        add_page_numbers, page_number_position, page_number_format, start_number,
    )
    asyncio.create_task(_schedule_job_cleanup(job_id))
    return JSONResponse(status_code=202, content={"job_id": job_id, "status": JobStatus.PENDING})


@router.post("/api/convert/rotate-pdf")
async def rotate_pdf(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    rotation: int = Form(90),
    pages: str = Form("all"),
    rotations: Optional[str] = Form(None),
):
    """Rotate PDF pages. rotation = 90/180/270 or per-page JSON map 'rotations'."""
    if not file.filename.lower().endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Please upload a PDF file.")
    if not rotations and rotation not in (90, 180, 270):
        raise HTTPException(status_code=400, detail="rotation must be 90, 180, or 270.")

    stem = Path(file.filename).stem
    tmp_dir = tempfile.mkdtemp()
    in_path = os.path.join(tmp_dir, "input.pdf")
    CHUNK = 1024 * 1024
    try:
        with open(in_path, "wb") as out_f:
            while True:
                chunk = await file.read(CHUNK)
                if not chunk:
                    break
                out_f.write(chunk)
    except Exception as e:
        shutil.rmtree(tmp_dir, ignore_errors=True)
        raise HTTPException(status_code=500, detail=f"Upload failed: {e}")
    with open(in_path, "rb") as f:
        if f.read(4) != b'%PDF':
            shutil.rmtree(tmp_dir, ignore_errors=True)
            raise HTTPException(status_code=400, detail="Not a valid PDF file.")

    job_id = str(uuid.uuid4())
    _job_store[job_id] = {
        "status": JobStatus.PENDING, "result_path": None, "tmp_dir": tmp_dir,
        "filename": f"{stem}_rotated.pdf", "content_type": "application/pdf", "error": None,
    }
    background_tasks.add_task(
        asyncio.get_event_loop().run_in_executor, _executor,
        _run_rotate_pdf_sync, job_id, in_path, tmp_dir, stem, rotation, pages, rotations,
    )
    asyncio.create_task(_schedule_job_cleanup(job_id))
    return JSONResponse(status_code=202, content={"job_id": job_id, "status": JobStatus.PENDING})


@router.post("/api/convert/add-watermark")
async def add_watermark(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    watermark_type: str = Form("text"),
    text: str = Form("CONFIDENTIAL"),
    image_file: UploadFile = File(None),
    font_family: str = Form("Helvetica"),
    font_size: int = Form(36),
    bold: bool = Form(False),
    italic: bool = Form(False),
    color: str = Form("#FF0000"),
    position: str = Form("center"),
    opacity: float = Form(0.5),
    rotation: int = Form(45),
    from_page: int = Form(1),
    to_page: int = Form(0),
    layer: str = Form("over"),
    scale: float = Form(1.0),
):
    """Add a text or image watermark to PDF pages with customizable parameters."""
    if not file.filename.lower().endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Please upload a PDF file.")

    stem = Path(file.filename).stem
    tmp_dir = tempfile.mkdtemp()
    in_path = os.path.join(tmp_dir, "input.pdf")
    CHUNK = 1024 * 1024
    try:
        with open(in_path, "wb") as out_f:
            while True:
                chunk = await file.read(CHUNK)
                if not chunk:
                    break
                out_f.write(chunk)
    except Exception as e:
        shutil.rmtree(tmp_dir, ignore_errors=True)
        raise HTTPException(status_code=500, detail=f"Upload failed: {e}")
    with open(in_path, "rb") as f:
        if f.read(4) != b'%PDF':
            shutil.rmtree(tmp_dir, ignore_errors=True)
            raise HTTPException(status_code=400, detail="Not a valid PDF file.")

    image_path = None
    if watermark_type == "image":
        if not image_file or not image_file.filename:
            shutil.rmtree(tmp_dir, ignore_errors=True)
            raise HTTPException(status_code=400, detail="An image file is required for image watermark.")
        ext = os.path.splitext(image_file.filename)[1] or ".png"
        image_path = os.path.join(tmp_dir, f"watermark_img{ext}")
        try:
            with open(image_path, "wb") as img_f:
                while True:
                    chunk = await image_file.read(CHUNK)
                    if not chunk:
                        break
                    img_f.write(chunk)
        except Exception as e:
            shutil.rmtree(tmp_dir, ignore_errors=True)
            raise HTTPException(status_code=500, detail=f"Image upload failed: {e}")

    job_id = str(uuid.uuid4())
    _job_store[job_id] = {
        "status": JobStatus.PENDING, "result_path": None, "tmp_dir": tmp_dir,
        "filename": f"{stem}_watermarked.pdf", "content_type": "application/pdf", "error": None,
    }
    _executor.submit(
        _run_add_watermark_sync, job_id, in_path, tmp_dir, stem,
        watermark_type, text, image_path, font_family, font_size,
        bold, italic, color, position, opacity, rotation, from_page, to_page, layer, scale,
    )
    asyncio.create_task(_schedule_job_cleanup(job_id))
    return JSONResponse(status_code=202, content={"job_id": job_id, "status": JobStatus.PENDING})


@router.post("/api/convert/add-page-numbers")
async def add_page_numbers(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    position: str = Form("bottom"),
    start: int = Form(1),
):
    """Add page numbers to each page. position = 'top' or 'bottom'. start = starting number."""
    if not file.filename.lower().endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Please upload a PDF file.")

    stem = Path(file.filename).stem
    tmp_dir = tempfile.mkdtemp()
    in_path = os.path.join(tmp_dir, "input.pdf")
    CHUNK = 1024 * 1024
    try:
        with open(in_path, "wb") as out_f:
            while True:
                chunk = await file.read(CHUNK)
                if not chunk:
                    break
                out_f.write(chunk)
    except Exception as e:
        shutil.rmtree(tmp_dir, ignore_errors=True)
        raise HTTPException(status_code=500, detail=f"Upload failed: {e}")
    with open(in_path, "rb") as f:
        if f.read(4) != b'%PDF':
            shutil.rmtree(tmp_dir, ignore_errors=True)
            raise HTTPException(status_code=400, detail="Not a valid PDF file.")

    job_id = str(uuid.uuid4())
    _job_store[job_id] = {
        "status": JobStatus.PENDING, "result_path": None, "tmp_dir": tmp_dir,
        "filename": f"{stem}_numbered.pdf", "content_type": "application/pdf", "error": None,
    }
    background_tasks.add_task(
        asyncio.get_event_loop().run_in_executor, _executor,
        _run_add_page_numbers_sync, job_id, in_path, tmp_dir, stem, position, start,
    )
    asyncio.create_task(_schedule_job_cleanup(job_id))
    return JSONResponse(status_code=202, content={"job_id": job_id, "status": JobStatus.PENDING})





@router.post("/api/convert/repair-pdf")
async def repair_pdf(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
):
    """Submit a PDF repair job. Returns job_id immediately."""
    if not file.filename.lower().endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Please upload a PDF file.")

    stem = Path(file.filename).stem
    tmp_dir = tempfile.mkdtemp()
    in_path = os.path.join(tmp_dir, "input.pdf")
    CHUNK = 1024 * 1024
    try:
        with open(in_path, "wb") as out_f:
            while True:
                chunk = await file.read(CHUNK)
                if not chunk:
                    break
                out_f.write(chunk)
    except Exception as e:
        shutil.rmtree(tmp_dir, ignore_errors=True)
        raise HTTPException(status_code=500, detail=f"Upload failed: {e}")

    job_id = str(uuid.uuid4())
    _job_store[job_id] = {
        "status": JobStatus.PENDING, "result_path": None, "tmp_dir": tmp_dir,
        "filename": f"{stem}_repaired.pdf", "content_type": "application/pdf", "error": None,
    }
    background_tasks.add_task(
        asyncio.get_event_loop().run_in_executor, _executor,
        _run_repair_pdf_sync, job_id, in_path, tmp_dir, stem,
    )
    asyncio.create_task(_schedule_job_cleanup(job_id))
    return JSONResponse(status_code=202, content={"job_id": job_id, "status": JobStatus.PENDING})


@router.post("/api/convert/pdf-to-pptx")
async def convert_pdf_to_pptx(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
):
    """Submit a PDF to PPTX conversion job."""
    if not file.filename.lower().endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Please upload a PDF file.")

    stem = Path(file.filename).stem
    tmp_dir = tempfile.mkdtemp()
    in_path = os.path.join(tmp_dir, "input.pdf")
    CHUNK = 1024 * 1024
    try:
        with open(in_path, "wb") as out_f:
            while True:
                chunk = await file.read(CHUNK)
                if not chunk:
                    break
                out_f.write(chunk)
    except Exception as e:
        shutil.rmtree(tmp_dir, ignore_errors=True)
        raise HTTPException(status_code=500, detail=f"Upload failed: {e}")

    job_id = str(uuid.uuid4())
    _job_store[job_id] = {
        "status": JobStatus.PENDING, "result_path": None, "tmp_dir": tmp_dir,
        "filename": f"{stem}.pptx", "content_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation", "error": None,
    }
    background_tasks.add_task(
        asyncio.get_event_loop().run_in_executor, _executor,
        _run_pdf_to_pptx_sync, job_id, in_path, tmp_dir, stem,
    )
    asyncio.create_task(_schedule_job_cleanup(job_id))
    return JSONResponse(status_code=202, content={"job_id": job_id, "status": JobStatus.PENDING})


@router.post("/api/convert/pdf-to-excel")
async def convert_pdf_to_excel(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
):
    """Submit a PDF to Excel conversion job."""
    if not file.filename.lower().endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Please upload a PDF file.")

    stem = Path(file.filename).stem
    tmp_dir = tempfile.mkdtemp()
    in_path = os.path.join(tmp_dir, "input.pdf")
    CHUNK = 1024 * 1024
    try:
        with open(in_path, "wb") as out_f:
            while True:
                chunk = await file.read(CHUNK)
                if not chunk:
                    break
                out_f.write(chunk)
    except Exception as e:
        shutil.rmtree(tmp_dir, ignore_errors=True)
        raise HTTPException(status_code=500, detail=f"Upload failed: {e}")

    job_id = str(uuid.uuid4())
    _job_store[job_id] = {
        "status": JobStatus.PENDING, "result_path": None, "tmp_dir": tmp_dir,
        "filename": f"{stem}.xlsx", "content_type": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "error": None,
    }
    background_tasks.add_task(
        asyncio.get_event_loop().run_in_executor, _executor,
        _run_pdf_to_excel_sync, job_id, in_path, tmp_dir, stem,
    )
    asyncio.create_task(_schedule_job_cleanup(job_id))
    return JSONResponse(status_code=202, content={"job_id": job_id, "status": JobStatus.PENDING})


@router.post("/api/convert/pptx-to-pdf")
async def convert_pptx_to_pdf(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
):
    """Submit a PowerPoint to PDF conversion job via LibreOffice."""
    soffice = _find_libreoffice()
    if not soffice:
        raise HTTPException(status_code=503, detail="LibreOffice is not installed.")

    ext = Path(file.filename).suffix.lower()
    if ext not in (".pptx", ".ppt", ".odp"):
        raise HTTPException(status_code=400, detail="Please upload a PowerPoint presentation (.pptx, .ppt, or .odp).")

    stem = Path(file.filename).stem
    tmp_dir = tempfile.mkdtemp()
    in_path = os.path.join(tmp_dir, f"{stem}{ext}")
    CHUNK = 1024 * 1024
    try:
        with open(in_path, "wb") as out_f:
            while True:
                chunk = await file.read(CHUNK)
                if not chunk:
                    break
                out_f.write(chunk)
    except Exception as e:
        shutil.rmtree(tmp_dir, ignore_errors=True)
        raise HTTPException(status_code=500, detail=f"Upload failed: {e}")

    job_id = str(uuid.uuid4())
    _job_store[job_id] = {
        "status": JobStatus.PENDING, "result_path": None, "tmp_dir": tmp_dir,
        "filename": f"{stem}.pdf", "content_type": "application/pdf", "error": None,
    }
    background_tasks.add_task(
        asyncio.get_event_loop().run_in_executor, _executor,
        _run_office_to_pdf_sync, job_id, in_path, tmp_dir, soffice,
    )
    asyncio.create_task(_schedule_job_cleanup(job_id))
    return JSONResponse(status_code=202, content={"job_id": job_id, "status": JobStatus.PENDING})


@router.post("/api/convert/excel-to-pdf")
async def convert_excel_to_pdf(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
):
    """Submit an Excel to PDF conversion job via LibreOffice."""
    soffice = _find_libreoffice()
    if not soffice:
        raise HTTPException(status_code=503, detail="LibreOffice is not installed.")

    ext = Path(file.filename).suffix.lower()
    if ext not in (".xlsx", ".xls", ".ods", ".csv"):
        raise HTTPException(status_code=400, detail="Please upload an Excel document (.xlsx, .xls, .ods, or .csv).")

    stem = Path(file.filename).stem
    tmp_dir = tempfile.mkdtemp()
    in_path = os.path.join(tmp_dir, f"{stem}{ext}")
    CHUNK = 1024 * 1024
    try:
        with open(in_path, "wb") as out_f:
            while True:
                chunk = await file.read(CHUNK)
                if not chunk:
                    break
                out_f.write(chunk)
    except Exception as e:
        shutil.rmtree(tmp_dir, ignore_errors=True)
        raise HTTPException(status_code=500, detail=f"Upload failed: {e}")

    job_id = str(uuid.uuid4())
    _job_store[job_id] = {
        "status": JobStatus.PENDING, "result_path": None, "tmp_dir": tmp_dir,
        "filename": f"{stem}.pdf", "content_type": "application/pdf", "error": None,
    }
    background_tasks.add_task(
        asyncio.get_event_loop().run_in_executor, _executor,
        _run_office_to_pdf_sync, job_id, in_path, tmp_dir, soffice,
    )
    asyncio.create_task(_schedule_job_cleanup(job_id))
    return JSONResponse(status_code=202, content={"job_id": job_id, "status": JobStatus.PENDING})


@router.post("/api/convert/flatten-pdf")
async def flatten_pdf(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
):
    """Submit a PDF form flatten job."""
    if not file.filename.lower().endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Please upload a PDF file.")

    stem = Path(file.filename).stem
    tmp_dir = tempfile.mkdtemp()
    in_path = os.path.join(tmp_dir, "input.pdf")
    CHUNK = 1024 * 1024
    try:
        with open(in_path, "wb") as out_f:
            while True:
                chunk = await file.read(CHUNK)
                if not chunk:
                    break
                out_f.write(chunk)
    except Exception as e:
        shutil.rmtree(tmp_dir, ignore_errors=True)
        raise HTTPException(status_code=500, detail=f"Upload failed: {e}")

    job_id = str(uuid.uuid4())
    _job_store[job_id] = {
        "status": JobStatus.PENDING, "result_path": None, "tmp_dir": tmp_dir,
        "filename": f"{stem}_flat.pdf", "content_type": "application/pdf", "error": None,
    }
    background_tasks.add_task(
        asyncio.get_event_loop().run_in_executor, _executor,
        _run_flatten_pdf_sync, job_id, in_path, tmp_dir, stem,
    )
    asyncio.create_task(_schedule_job_cleanup(job_id))
    return JSONResponse(status_code=202, content={"job_id": job_id, "status": JobStatus.PENDING})


@router.post("/api/convert/lock-pdf")
async def lock_pdf(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    password: str = Form(...),
    lock_mode: str = Form("open"),
    allow_print: bool = Form(True),
    allow_copy: bool = Form(True),
):
    """Encrypt a PDF document with a user-specified password using 256-bit AES."""
    if not file.filename.lower().endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Please upload a PDF file.")
    if not password or len(password.strip()) == 0:
        raise HTTPException(status_code=400, detail="A password is required to lock the PDF.")

    stem = Path(file.filename).stem
    tmp_dir = tempfile.mkdtemp()
    in_path = os.path.join(tmp_dir, "input.pdf")
    CHUNK = 1024 * 1024
    try:
        with open(in_path, "wb") as out_f:
            while True:
                chunk = await file.read(CHUNK)
                if not chunk:
                    break
                out_f.write(chunk)
    except Exception as e:
        shutil.rmtree(tmp_dir, ignore_errors=True)
        raise HTTPException(status_code=500, detail=f"Upload failed: {e}")

    job_id = str(uuid.uuid4())
    _job_store[job_id] = {
        "status": JobStatus.PENDING, "result_path": None, "tmp_dir": tmp_dir,
        "filename": f"{stem}_protected.pdf", "content_type": "application/pdf", "error": None,
    }
    background_tasks.add_task(
        asyncio.get_event_loop().run_in_executor, _executor,
        _run_lock_pdf_sync, job_id, in_path, tmp_dir, stem, password, lock_mode, allow_print, allow_copy,
    )
    asyncio.create_task(_schedule_job_cleanup(job_id))
    return JSONResponse(status_code=202, content={"job_id": job_id, "status": JobStatus.PENDING})


@router.post("/api/convert/check-pdf-lock")
async def check_pdf_lock(
    file: UploadFile = File(...),
):
    """Smart Lock Detector: Determines if a PDF can be unlocked automatically without a password."""
    if not file.filename or not file.filename.lower().endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Uploaded file must be a PDF.")

    try:
        import pikepdf, io
        pdf_bytes = await file.read()
        
        try:
            pdf = pikepdf.Pdf.open(io.BytesIO(pdf_bytes), password="")
            is_enc = pdf.is_encrypted
            pdf.close()
            if not is_enc:
                return JSONResponse(content={"status": "unencrypted", "can_auto_unlock": True, "message": "This PDF is not password-protected."})
            return JSONResponse(content={"status": "auto_unlockable", "can_auto_unlock": True, "message": "Instant unlock available! No password required."})
        except pikepdf.PasswordError:
            return JSONResponse(content={"status": "requires_password", "can_auto_unlock": False, "message": "Strict Open Password required to unlock."})
        except Exception as e:
            return JSONResponse(content={"status": "requires_password", "can_auto_unlock": False, "message": str(e)})
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Check failed: {e}")


@router.post("/api/convert/unlock-pdf")
async def unlock_pdf(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    password: str = Form(""),
):
    """Remove password protection from a PDF document."""
    if not file.filename.lower().endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Please upload a PDF file.")

    stem = Path(file.filename).stem
    tmp_dir = tempfile.mkdtemp()
    in_path = os.path.join(tmp_dir, "input.pdf")
    CHUNK = 1024 * 1024
    try:
        with open(in_path, "wb") as out_f:
            while True:
                chunk = await file.read(CHUNK)
                if not chunk:
                    break
                out_f.write(chunk)
    except Exception as e:
        shutil.rmtree(tmp_dir, ignore_errors=True)
        raise HTTPException(status_code=500, detail=f"Upload failed: {e}")

    job_id = str(uuid.uuid4())
    _job_store[job_id] = {
        "status": JobStatus.PENDING, "result_path": None, "tmp_dir": tmp_dir,
        "filename": f"{stem}_unlocked.pdf", "content_type": "application/pdf", "error": None,
    }
    background_tasks.add_task(
        asyncio.get_event_loop().run_in_executor, _executor,
        _run_unlock_pdf_sync, job_id, in_path, tmp_dir, stem, password,
    )
    asyncio.create_task(_schedule_job_cleanup(job_id))
    return JSONResponse(status_code=202, content={"job_id": job_id, "status": JobStatus.PENDING})


@router.post("/api/convert/redact-pdf")
async def redact_pdf(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    redactions: str = Form(...),
):
    """Apply permanent redactions (rectangles + text search) to a PDF document.

    The 'redactions' form field must be a JSON string with shape:
    { "items": [
        { "page": 0, "type": "rect", "rect": [x0, y0, x1, y1], "color": "#000000" },
        { "page": 0, "type": "text", "text": "sensitive phrase", "color": "#000000" }
    ]}
    Coordinates are in PDF user-space points (origin bottom-left per PDF spec).
    """
    if not file.filename.lower().endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Please upload a PDF file.")

    stem = Path(file.filename).stem
    tmp_dir = tempfile.mkdtemp()
    in_path = os.path.join(tmp_dir, "input.pdf")
    CHUNK = 1024 * 1024
    try:
        with open(in_path, "wb") as out_f:
            while True:
                chunk = await file.read(CHUNK)
                if not chunk:
                    break
                out_f.write(chunk)
    except Exception as e:
        shutil.rmtree(tmp_dir, ignore_errors=True)
        raise HTTPException(status_code=500, detail=f"Upload failed: {e}")

    job_id = str(uuid.uuid4())
    _job_store[job_id] = {
        "status": JobStatus.PENDING, "result_path": None, "tmp_dir": tmp_dir,
        "filename": f"{stem}_redacted.pdf", "content_type": "application/pdf", "error": None,
    }
    background_tasks.add_task(
        asyncio.get_event_loop().run_in_executor, _executor,
        _run_redact_pdf_sync, job_id, in_path, tmp_dir, stem, redactions,
    )
    asyncio.create_task(_schedule_job_cleanup(job_id))
    return JSONResponse(status_code=202, content={"job_id": job_id, "status": JobStatus.PENDING})


@router.post("/api/convert/html-to-pdf")
async def convert_html_to_pdf(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(None),
    html: str = Form(""),
):
    """Convert uploaded HTML file or raw HTML text to PDF."""
    html_content = ""
    stem = "document"

    if file and file.filename:
        stem = Path(file.filename).stem
        content_bytes = await file.read()
        html_content = content_bytes.decode("utf-8", errors="ignore")
    elif html.strip():
        html_content = html
        stem = "webpage"
    else:
        raise HTTPException(status_code=400, detail="Please upload an HTML file or enter HTML content.")

    tmp_dir = tempfile.mkdtemp()
    job_id = str(uuid.uuid4())
    _job_store[job_id] = {
        "status": JobStatus.PENDING, "result_path": None, "tmp_dir": tmp_dir,
        "filename": f"{stem}.pdf", "content_type": "application/pdf", "error": None,
    }
    background_tasks.add_task(
        asyncio.get_event_loop().run_in_executor, _executor,
        _run_html_to_pdf_sync, job_id, html_content, tmp_dir, stem,
    )
    asyncio.create_task(_schedule_job_cleanup(job_id))
    return JSONResponse(status_code=202, content={"job_id": job_id, "status": JobStatus.PENDING})


@router.post("/api/convert/pdf-to-pdfa")
async def convert_pdf_to_pdfa(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    conformance: str = Form("PDF/A-2b"),
    allow_downgrade: bool = Form(True),
):
    """Submit a PDF to PDF/A compliance job."""
    if not file.filename.lower().endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Please upload a PDF file.")

    stem = Path(file.filename).stem
    tmp_dir = tempfile.mkdtemp()
    in_path = os.path.join(tmp_dir, "input.pdf")
    CHUNK = 1024 * 1024
    try:
        with open(in_path, "wb") as out_f:
            while True:
                chunk = await file.read(CHUNK)
                if not chunk:
                    break
                out_f.write(chunk)
    except Exception as e:
        shutil.rmtree(tmp_dir, ignore_errors=True)
        raise HTTPException(status_code=500, detail=f"Upload failed: {e}")

    job_id = str(uuid.uuid4())
    _job_store[job_id] = {
        "status": JobStatus.PENDING, "result_path": None, "tmp_dir": tmp_dir,
        "filename": f"{stem}_pdfa.pdf", "content_type": "application/pdf", "error": None,
    }
    background_tasks.add_task(
        asyncio.get_event_loop().run_in_executor, _executor,
        _run_pdf_to_pdfa_sync, job_id, in_path, tmp_dir, stem, conformance, allow_downgrade,
    )
    asyncio.create_task(_schedule_job_cleanup(job_id))
    return JSONResponse(status_code=202, content={"job_id": job_id, "status": JobStatus.PENDING})


@router.post("/api/convert/ocr-to-word")
async def convert_ocr_to_word(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    language: str = Form("eng"),
):
    """Perform OCR on an uploaded image or PDF and output a Word (.docx) document."""
    ext = Path(file.filename).suffix.lower()
    allowed_exts = ('.pdf', '.png', '.jpg', '.jpeg', '.webp', '.tiff', '.tif', '.bmp')
    if ext not in allowed_exts:
        raise HTTPException(status_code=400, detail=f"Unsupported file format: {ext}. Allowed: PDF, PNG, JPG, WEBP, TIFF.")

    stem = Path(file.filename).stem
    tmp_dir = tempfile.mkdtemp()
    in_path = os.path.join(tmp_dir, f"input{ext}")
    CHUNK = 1024 * 1024
    try:
        with open(in_path, "wb") as out_f:
            while True:
                chunk = await file.read(CHUNK)
                if not chunk:
                    break
                out_f.write(chunk)
    except Exception as e:
        shutil.rmtree(tmp_dir, ignore_errors=True)
        raise HTTPException(status_code=500, detail=f"Upload failed: {e}")

    job_id = str(uuid.uuid4())
    _job_store[job_id] = {
        "status": JobStatus.PENDING, "result_path": None, "tmp_dir": tmp_dir,
        "filename": f"{stem}_ocr.docx",
        "content_type": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "error": None,
    }
    background_tasks.add_task(
        asyncio.get_event_loop().run_in_executor, _executor,
        _run_ocr_to_word_sync, job_id, in_path, tmp_dir, stem, language,
    )
    asyncio.create_task(_schedule_job_cleanup(job_id))
    return JSONResponse(status_code=202, content={"job_id": job_id, "status": JobStatus.PENDING})


@router.post("/api/convert/pdf-to-markdown")
async def convert_pdf_to_markdown(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
):
    """Convert PDF document to Markdown (.md) format with headers, paragraphs, lists, and tables."""
    if not file.filename.lower().endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Please upload a PDF file.")

    stem = Path(file.filename).stem
    tmp_dir = tempfile.mkdtemp()
    in_path = os.path.join(tmp_dir, "input.pdf")
    CHUNK = 1024 * 1024
    try:
        with open(in_path, "wb") as out_f:
            while True:
                chunk = await file.read(CHUNK)
                if not chunk:
                    break
                out_f.write(chunk)
    except Exception as e:
        shutil.rmtree(tmp_dir, ignore_errors=True)
        raise HTTPException(status_code=500, detail=f"Upload failed: {e}")

    with open(in_path, "rb") as f:
        if f.read(4) != b'%PDF':
            shutil.rmtree(tmp_dir, ignore_errors=True)
            raise HTTPException(status_code=400, detail="Not a valid PDF file.")

    job_id = str(uuid.uuid4())
    _job_store[job_id] = {
        "status": JobStatus.PENDING, "result_path": None, "tmp_dir": tmp_dir,
        "filename": f"{stem}.md", "content_type": "text/markdown", "error": None,
    }
    background_tasks.add_task(
        asyncio.get_event_loop().run_in_executor, _executor,
        _run_pdf_to_markdown_sync, job_id, in_path, tmp_dir, stem,
    )
    asyncio.create_task(_schedule_job_cleanup(job_id))
    return JSONResponse(status_code=202, content={"job_id": job_id, "status": JobStatus.PENDING})


@router.post("/api/convert/compare-pdf")
async def compare_pdf(
    background_tasks: BackgroundTasks,
    files: List[UploadFile] = File(...),
):
    """Submit a side-by-side PDF comparison job (requires exactly 2 PDF files)."""
    if len(files) != 2:
        raise HTTPException(status_code=400, detail="Please upload exactly 2 PDF files to compare.")

    for f in files:
        if not f.filename.lower().endswith('.pdf'):
            raise HTTPException(status_code=400, detail=f"'{f.filename}' is not a PDF file.")

    tmp_dir = tempfile.mkdtemp()
    file_paths = []
    CHUNK = 1024 * 1024
    for f in files:
        in_path = os.path.join(tmp_dir, Path(f.filename).name)
        try:
            with open(in_path, "wb") as out_f:
                while True:
                    chunk = await f.read(CHUNK)
                    if not chunk:
                        break
                    out_f.write(chunk)
        except Exception as e:
            shutil.rmtree(tmp_dir, ignore_errors=True)
            raise HTTPException(status_code=500, detail=f"Upload failed: {e}")
        file_paths.append(in_path)

    stem1 = Path(files[0].filename).stem
    stem2 = Path(files[1].filename).stem

    job_id = str(uuid.uuid4())
    _job_store[job_id] = {
        "status": JobStatus.PENDING, "result_path": None, "tmp_dir": tmp_dir,
        "filename": f"{stem1}_vs_{stem2}_comparison.pdf", "content_type": "application/pdf", "error": None,
        "comparison_data": None,
    }
    background_tasks.add_task(
        asyncio.get_event_loop().run_in_executor, _executor,
        _run_compare_pdf_sync, job_id, file_paths[0], file_paths[1], tmp_dir, stem1, stem2,
    )
    asyncio.create_task(_schedule_job_cleanup(job_id))
    return JSONResponse(status_code=202, content={"job_id": job_id, "status": JobStatus.PENDING})


# ─── Text to PDF & PDF to HTML Conversion Handlers ──────────────────────────

def _run_text_to_pdf_sync(job_id: str, in_path: str, tmp_dir: str, stem: str):
    """Blocking Text to PDF conversion using ReportLab."""
    try:
        import html
        from reportlab.lib.pagesizes import letter
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib import colors

        _job_store[job_id]["status"] = JobStatus.PROCESSING
        pdf_path = os.path.join(tmp_dir, f"{stem}.pdf")

        with open(in_path, 'r', encoding='utf-8', errors='replace') as f:
            content = f.read()

        doc = SimpleDocTemplate(
            pdf_path,
            pagesize=letter,
            rightMargin=54, leftMargin=54,
            topMargin=54, bottomMargin=54
        )

        styles = getSampleStyleSheet()
        normal_style = ParagraphStyle(
            'CustomText',
            parent=styles['Normal'],
            fontName='Helvetica',
            fontSize=10,
            leading=14,
            textColor=colors.HexColor('#0F172A')
        )

        story = []
        paragraphs = content.split('\n\n')
        for p_text in paragraphs:
            if p_text.strip():
                escaped = html.escape(p_text.strip()).replace('\n', '<br />')
                story.append(Paragraph(escaped, normal_style))
                story.append(Spacer(1, 10))

        if not story:
            story.append(Paragraph("", normal_style))

        doc.build(story)

        if not os.path.isfile(pdf_path):
            _job_store[job_id]["status"] = JobStatus.ERROR
            _job_store[job_id]["error"] = "Failed to generate PDF document."
            return

        _job_store[job_id]["status"] = JobStatus.DONE
        _job_store[job_id]["result_path"] = pdf_path
    except Exception as e:
        _job_store[job_id]["status"] = JobStatus.ERROR
        _job_store[job_id]["error"] = str(e)


def _run_pdf_to_html_sync(job_id: str, in_path: str, tmp_dir: str, stem: str):
    """Blocking PDF to HTML conversion using PyMuPDF (fitz) with responsive web page layout."""
    try:
        import fitz
        import base64
        import re

        _job_store[job_id]["status"] = JobStatus.PROCESSING
        html_path = os.path.join(tmp_dir, f"{stem}.html")

        doc = fitz.open(in_path)
        html_chunks = [
            '<!DOCTYPE html>',
            '<html lang="en">',
            '<head>',
            '  <meta charset="utf-8">',
            '  <meta name="viewport" content="width=device-width, initial-scale=1.0">',
            '  <title>Converted Web Page</title>',
            '  <style>',
            '    * { box-sizing: border-box; }',
            '    html, body {',
            '      width: 100%;',
            '      min-height: 100vh;',
            '      margin: 0;',
            '      padding: 0;',
            '      background-color: #ffffff;',
            '      color: #0f172a;',
            '      font-family: system-ui, -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, "Helvetica Neue", Arial, sans-serif;',
            '      line-height: 1.6;',
            '    }',
            '    .web-page-wrapper {',
            '      width: 100%;',
            '      max-width: 1200px;',
            '      margin: 0 auto;',
            '      padding: 2rem 1.5rem;',
            '    }',
            '    .web-page-section {',
            '      position: relative;',
            '      width: 100%;',
            '      margin-bottom: 2rem;',
            '    }',
            '    p { margin: 0.4em 0; }',
            '    img { max-width: 100%; height: auto; display: inline-block; }',
            '  </style>',
            '</head>',
            '<body>',
            '  <div class="web-page-wrapper">'
        ]

        for i, page in enumerate(doc):
            page_html = page.get_text("html")

            # Clean up raw PDF font names & serif fallbacks
            page_html = re.sub(r'font-family:[^;"]*', 'font-family: system-ui, -apple-system, sans-serif', page_html)

            # Strip fixed container width/height styles from PyMuPDF page div
            page_html = re.sub(r'style="[^"]*width:\s*\d+\.?\d*pt;[^"]*"', 'style="width:100%; position:relative;"', page_html)

            # Embed page images as base64 data URIs
            try:
                seen_xrefs = set()
                img_tags = []
                for img_info in page.get_images(full=True):
                    xref = img_info[0]
                    if xref in seen_xrefs:
                        continue
                    seen_xrefs.add(xref)
                    base_img = doc.extract_image(xref)
                    if base_img and "image" in base_img:
                        ext = base_img.get("ext", "png")
                        b64 = base64.b64encode(base_img["image"]).decode("utf-8")
                        img_tags.append(f'<div style="margin: 1rem 0;"><img src="data:image/{ext};base64,{b64}" alt="Embedded Image" /></div>')
                if img_tags and "<img" not in page_html.lower():
                    page_html += "\n" + "\n".join(img_tags)
            except Exception:
                pass

            html_chunks.append(f'    <div class="web-page-section" id="section-{i+1}">')
            html_chunks.append(page_html)
            html_chunks.append('    </div>')

        html_chunks.append('  </div>')
        html_chunks.append('</body></html>')
        doc.close()

        with open(html_path, 'w', encoding='utf-8') as f:
            f.write("\n".join(html_chunks))

        if not os.path.isfile(html_path):
            _job_store[job_id]["status"] = JobStatus.ERROR
            _job_store[job_id]["error"] = "Failed to generate HTML document."
            return

        _job_store[job_id]["status"] = JobStatus.DONE
        _job_store[job_id]["result_path"] = html_path
    except Exception as e:
        _job_store[job_id]["status"] = JobStatus.ERROR
        _job_store[job_id]["error"] = str(e)


@router.post("/api/convert/text-to-pdf")
async def text_to_pdf(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
):
    """Submit an async Text to PDF conversion job."""
    if not file.filename.lower().endswith(('.txt', '.text', '.log', '.md', '.csv')):
        raise HTTPException(status_code=400, detail="Uploaded file must be a text document (.txt, .log, .md, .csv).")

    stem = Path(file.filename).stem
    tmp_dir = tempfile.mkdtemp()
    in_path = os.path.join(tmp_dir, Path(file.filename).name)

    CHUNK = 1024 * 1024
    try:
        with open(in_path, "wb") as out_f:
            while True:
                chunk = await file.read(CHUNK)
                if not chunk:
                    break
                out_f.write(chunk)
    except Exception as e:
        shutil.rmtree(tmp_dir, ignore_errors=True)
        raise HTTPException(status_code=500, detail=f"Upload failed: {e}")

    job_id = str(uuid.uuid4())
    _job_store[job_id] = {
        "status": JobStatus.PENDING, "result_path": None, "tmp_dir": tmp_dir,
        "filename": f"{stem}.pdf", "content_type": "application/pdf", "error": None,
    }
    background_tasks.add_task(
        asyncio.get_event_loop().run_in_executor, _executor,
        _run_text_to_pdf_sync, job_id, in_path, tmp_dir, stem,
    )
    asyncio.create_task(_schedule_job_cleanup(job_id))
    return JSONResponse(status_code=202, content={"job_id": job_id, "status": JobStatus.PENDING})


@router.post("/api/convert/pdf-to-html")
async def pdf_to_html(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
):
    """Submit an async PDF to HTML conversion job."""
    if not file.filename.lower().endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Uploaded file must be a PDF document.")

    stem = Path(file.filename).stem
    tmp_dir = tempfile.mkdtemp()
    in_path = os.path.join(tmp_dir, Path(file.filename).name)

    CHUNK = 1024 * 1024
    try:
        with open(in_path, "wb") as out_f:
            while True:
                chunk = await file.read(CHUNK)
                if not chunk:
                    break
                out_f.write(chunk)
    except Exception as e:
        shutil.rmtree(tmp_dir, ignore_errors=True)
        raise HTTPException(status_code=500, detail=f"Upload failed: {e}")

    job_id = str(uuid.uuid4())
    _job_store[job_id] = {
        "status": JobStatus.PENDING, "result_path": None, "tmp_dir": tmp_dir,
        "filename": f"{stem}.html", "content_type": "text/html", "error": None,
    }
    background_tasks.add_task(
        asyncio.get_event_loop().run_in_executor, _executor,
        _run_pdf_to_html_sync, job_id, in_path, tmp_dir, stem,
    )
    asyncio.create_task(_schedule_job_cleanup(job_id))
    return JSONResponse(status_code=202, content={"job_id": job_id, "status": JobStatus.PENDING})


# ─── Word <-> PPTX and Word <-> Excel Conversion Handlers ─────────────────

def _extract_images_from_docx_paragraph(p, doc, tmp_dir):
    img_paths = []
    try:
        blips = p._element.xpath('.//a:blip/@r:embed')
        for rId in blips:
            if rId in doc.part.related_parts:
                image_part = doc.part.related_parts[rId]
                image_bytes = image_part.blob
                ext = image_part.content_type.split('/')[-1]
                if ext == 'jpeg':
                    ext = 'jpg'
                img_path = os.path.join(tmp_dir, f"docx_img_{uuid.uuid4().hex}.{ext}")
                with open(img_path, 'wb') as img_f:
                    img_f.write(image_bytes)
                img_paths.append(img_path)
    except Exception:
        pass
    return img_paths


def _run_word_to_pptx_sync(job_id: str, in_path: str, tmp_dir: str, stem: str):
    """Blocking Word (.docx) to PowerPoint (.pptx) conversion."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from docx import Document
        from docx.text.paragraph import Paragraph
        from docx.table import Table

        _job_store[job_id]["status"] = JobStatus.PROCESSING
        out_path = os.path.join(tmp_dir, f"{stem}.pptx")

        doc = Document(in_path)
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]

        current_slide = None
        current_tf = None

        def create_new_slide(title_text=""):
            nonlocal current_slide, current_tf
            current_slide = prs.slides.add_slide(blank_slide_layout)

            txBox = current_slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.4), Inches(1.2))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title_text or stem
            p.font.size = Pt(28)
            p.font.bold = True

            content_box = current_slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(5.0))
            current_tf = content_box.text_frame
            current_tf.word_wrap = True

        for element in doc.element.body:
            if element.tag.endswith('p'):
                p = Paragraph(element, doc)
                p_text = p.text.strip()
                style_name = p.style.name if p.style else ""

                img_paths = _extract_images_from_docx_paragraph(p, doc, tmp_dir)
                for img_p in img_paths:
                    if current_slide is None:
                        create_new_slide(title_text=stem)
                    try:
                        current_slide.shapes.add_picture(img_p, Inches(0.8), Inches(2.0), width=Inches(4.5))
                    except Exception:
                        pass

                if not p_text:
                    continue

                if style_name.startswith("Heading 1") or (style_name.startswith("Heading") and current_slide is None):
                    create_new_slide(title_text=p_text)
                elif style_name.startswith("Heading"):
                    if current_slide is None:
                        create_new_slide(title_text=p_text)
                    else:
                        p_pptx = current_tf.add_paragraph() if current_tf.paragraphs[0].text else current_tf.paragraphs[0]
                        p_pptx.text = p_text
                        p_pptx.font.size = Pt(20)
                        p_pptx.font.bold = True
                else:
                    if current_slide is None:
                        create_new_slide(title_text=stem)
                    p_pptx = current_tf.add_paragraph() if current_tf.paragraphs[0].text else current_tf.paragraphs[0]
                    p_pptx.text = p_text
                    p_pptx.font.size = Pt(14)
                    if style_name.startswith("List") or p._element.xpath('.//w:numPr'):
                        p_pptx.level = 1

            elif element.tag.endswith('tbl'):
                tbl = Table(element, doc)
                if current_slide is None:
                    create_new_slide(title_text=stem)

                rows_count = len(tbl.rows)
                cols_count = len(tbl.columns) if rows_count > 0 else 0
                if rows_count > 0 and cols_count > 0:
                    try:
                        table_shape = current_slide.shapes.add_table(
                            rows_count, cols_count, Inches(0.8), Inches(2.0), Inches(8.4), Inches(0.4 * rows_count)
                        )
                        table = table_shape.table
                        for r_idx, row in enumerate(tbl.rows):
                            for c_idx, cell in enumerate(row.cells):
                                if c_idx < cols_count:
                                    table.cell(r_idx, c_idx).text = cell.text.strip()
                    except Exception:
                        pass

        if current_slide is None:
            create_new_slide(title_text=stem)

        prs.save(out_path)
        _job_store[job_id]["status"] = JobStatus.DONE
        _job_store[job_id]["result_path"] = out_path
    except Exception as e:
        _job_store[job_id]["status"] = JobStatus.ERROR
        _job_store[job_id]["error"] = str(e)


def _run_pptx_to_word_sync(job_id: str, in_path: str, tmp_dir: str, stem: str):
    """Blocking PowerPoint (.pptx) to Word (.docx) conversion."""
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from docx import Document
        from docx.shared import Inches

        _job_store[job_id]["status"] = JobStatus.PROCESSING
        out_path = os.path.join(tmp_dir, f"{stem}.docx")

        prs = Presentation(in_path)
        doc = Document()

        for s_idx, slide in enumerate(prs.slides, 1):
            title_text = ""
            if slide.shapes.title and slide.shapes.title.text:
                title_text = slide.shapes.title.text.strip()

            doc.add_heading(title_text or f"Slide {s_idx}", level=1)

            for shape in slide.shapes:
                if shape == slide.shapes.title:
                    continue

                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        text = p.text.strip()
                        if text:
                            if p.level > 0:
                                doc.add_paragraph(text, style='List Bullet')
                            else:
                                doc.add_paragraph(text)

                if shape.has_table:
                    table = shape.table
                    rows_count = len(table.rows)
                    cols_count = len(table.columns)
                    if rows_count > 0 and cols_count > 0:
                        docx_tbl = doc.add_table(rows=rows_count, cols=cols_count)
                        docx_tbl.style = 'Table Grid'
                        for r in range(rows_count):
                            for c in range(cols_count):
                                docx_tbl.cell(r, c).text = table.cell(r, c).text.strip()
                        doc.add_paragraph()

                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or hasattr(shape, "image"):
                    try:
                        image = shape.image
                        img_bytes = image.blob
                        ext = image.ext
                        img_path = os.path.join(tmp_dir, f"pptx_img_{uuid.uuid4().hex}.{ext}")
                        with open(img_path, 'wb') as img_f:
                            img_f.write(img_bytes)
                        doc.add_picture(img_path, width=Inches(4.5))
                    except Exception:
                        pass

        doc.save(out_path)
        _job_store[job_id]["status"] = JobStatus.DONE
        _job_store[job_id]["result_path"] = out_path
    except Exception as e:
        _job_store[job_id]["status"] = JobStatus.ERROR
        _job_store[job_id]["error"] = str(e)


def _run_word_to_excel_sync(job_id: str, in_path: str, tmp_dir: str, stem: str):
    """Blocking Word (.docx) to Excel (.xlsx) conversion."""
    try:
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
        from docx import Document

        _job_store[job_id]["status"] = JobStatus.PROCESSING
        out_path = os.path.join(tmp_dir, f"{stem}.xlsx")

        doc = Document(in_path)
        wb = openpyxl.Workbook()
        default_sheet = wb.active
        default_sheet.title = "Overview"

        header_font = Font(name='Calibri', size=11, bold=True, color='FFFFFF')
        header_fill = PatternFill(start_color='1F4E78', end_color='1F4E78', fill_type='solid')
        header_align = Alignment(horizontal='center', vertical='center', wrap_text=True)
        thin_border = Border(
            left=Side(style='thin', color='D9D9D9'),
            right=Side(style='thin', color='D9D9D9'),
            top=Side(style='thin', color='D9D9D9'),
            bottom=Side(style='thin', color='D9D9D9')
        )

        def style_header_row(ws, row_idx=1):
            for cell in ws[row_idx]:
                cell.font = header_font
                cell.fill = header_fill
                cell.alignment = header_align

        def auto_fit_columns(ws):
            for col in ws.columns:
                max_len = 0
                for cell in col:
                    cell.border = thin_border
                    val_str = str(cell.value or '')
                    if len(val_str) > max_len:
                        max_len = len(val_str)
                col_letter = openpyxl.utils.get_column_letter(col[0].column)
                ws.column_dimensions[col_letter].width = max(max_len + 4, 12)

        has_tables = len(doc.tables) > 0
        for idx, table in enumerate(doc.tables, 1):
            sheet_title = f"Table {idx}"
            if idx == 1 and default_sheet.title == "Overview":
                ws = default_sheet
                ws.title = sheet_title
            else:
                ws = wb.create_sheet(title=sheet_title)

            for r_idx, row in enumerate(table.rows, 1):
                row_data = [cell.text.strip() for cell in row.cells]
                ws.append(row_data)

            if ws.max_row >= 1:
                style_header_row(ws, 1)
            auto_fit_columns(ws)

        text_data = []
        for p in doc.paragraphs:
            p_text = p.text.strip()
            if not p_text:
                continue
            style_name = p.style.name if p.style else ""
            if style_name.startswith("Heading"):
                text_data.append(["Heading", p_text, ""])
            elif ":" in p_text and not p_text.startswith("http"):
                parts = p_text.split(":", 1)
                text_data.append(["Key-Value", parts[0].strip(), parts[1].strip()])
            elif style_name.startswith("List") or p._element.xpath('.//w:numPr'):
                text_data.append(["List Item", p_text, ""])
            else:
                text_data.append(["Paragraph", p_text, ""])

        if text_data:
            ws_summary = wb.create_sheet(title="Document Summary") if has_tables else default_sheet
            if not has_tables:
                ws_summary.title = "Document Summary"
            ws_summary.append(["Category", "Content / Key", "Details / Value"])
            for row in text_data:
                ws_summary.append(row)
            style_header_row(ws_summary, 1)
            auto_fit_columns(ws_summary)

        if "Overview" in wb.sheetnames and len(wb.sheetnames) > 1:
            ws_ov = wb["Overview"]
            if ws_ov.max_row <= 1 and ws_ov.cell(1, 1).value is None:
                wb.remove(ws_ov)

        wb.save(out_path)
        _job_store[job_id]["status"] = JobStatus.DONE
        _job_store[job_id]["result_path"] = out_path
    except Exception as e:
        _job_store[job_id]["status"] = JobStatus.ERROR
        _job_store[job_id]["error"] = str(e)


def _run_excel_to_word_sync(job_id: str, in_path: str, tmp_dir: str, stem: str):
    """Blocking Excel (.xlsx) to Word (.docx) conversion."""
    try:
        import openpyxl
        from docx import Document

        _job_store[job_id]["status"] = JobStatus.PROCESSING
        out_path = os.path.join(tmp_dir, f"{stem}.docx")

        wb = openpyxl.load_workbook(in_path, data_only=True)
        doc = Document()

        for sheet in wb.worksheets:
            rows = list(sheet.iter_rows(values_only=True))
            while rows and not any(cell is not None and str(cell).strip() != "" for cell in rows[-1]):
                rows.pop()
            if not rows:
                continue

            doc.add_heading(sheet.title, level=1)

            max_cols = max(len(r) for r in rows) if rows else 0
            if max_cols > 0:
                table = doc.add_table(rows=len(rows), cols=max_cols)
                table.style = 'Table Grid'
                for r_idx, row_values in enumerate(rows):
                    for c_idx in range(max_cols):
                        val = row_values[c_idx] if c_idx < len(row_values) else ""
                        cell_str = str(val).strip() if val is not None else ""
                        cell = table.cell(r_idx, c_idx)
                        cell.text = cell_str
                        if r_idx == 0:
                            for p in cell.paragraphs:
                                for run in p.runs:
                                    run.bold = True
                doc.add_paragraph()

        doc.save(out_path)
        _job_store[job_id]["status"] = JobStatus.DONE
        _job_store[job_id]["result_path"] = out_path
    except Exception as e:
        _job_store[job_id]["status"] = JobStatus.ERROR
        _job_store[job_id]["error"] = str(e)


@router.post("/api/convert/word-to-pptx")
async def convert_word_to_pptx(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
):
    """Submit an async Word (.docx) to PowerPoint (.pptx) conversion job."""
    if not file.filename.lower().endswith(('.docx', '.doc')):
        raise HTTPException(status_code=400, detail="Uploaded file must be a Word document (.docx).")

    stem = Path(file.filename).stem
    tmp_dir = tempfile.mkdtemp()
    in_path = os.path.join(tmp_dir, Path(file.filename).name)

    CHUNK = 1024 * 1024
    try:
        with open(in_path, "wb") as out_f:
            while True:
                chunk = await file.read(CHUNK)
                if not chunk:
                    break
                out_f.write(chunk)
    except Exception as e:
        shutil.rmtree(tmp_dir, ignore_errors=True)
        raise HTTPException(status_code=500, detail=f"Upload failed: {e}")

    job_id = str(uuid.uuid4())
    _job_store[job_id] = {
        "status": JobStatus.PENDING, "result_path": None, "tmp_dir": tmp_dir,
        "filename": f"{stem}.pptx",
        "content_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "error": None,
    }
    background_tasks.add_task(
        asyncio.get_event_loop().run_in_executor, _executor,
        _run_word_to_pptx_sync, job_id, in_path, tmp_dir, stem,
    )
    asyncio.create_task(_schedule_job_cleanup(job_id))
    return JSONResponse(status_code=202, content={"job_id": job_id, "status": JobStatus.PENDING})


@router.post("/api/convert/pptx-to-word")
async def convert_pptx_to_word(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
):
    """Submit an async PowerPoint (.pptx) to Word (.docx) conversion job."""
    if not file.filename.lower().endswith(('.pptx', '.ppt')):
        raise HTTPException(status_code=400, detail="Uploaded file must be a PowerPoint presentation (.pptx).")

    stem = Path(file.filename).stem
    tmp_dir = tempfile.mkdtemp()
    in_path = os.path.join(tmp_dir, Path(file.filename).name)

    CHUNK = 1024 * 1024
    try:
        with open(in_path, "wb") as out_f:
            while True:
                chunk = await file.read(CHUNK)
                if not chunk:
                    break
                out_f.write(chunk)
    except Exception as e:
        shutil.rmtree(tmp_dir, ignore_errors=True)
        raise HTTPException(status_code=500, detail=f"Upload failed: {e}")

    job_id = str(uuid.uuid4())
    _job_store[job_id] = {
        "status": JobStatus.PENDING, "result_path": None, "tmp_dir": tmp_dir,
        "filename": f"{stem}.docx",
        "content_type": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "error": None,
    }
    background_tasks.add_task(
        asyncio.get_event_loop().run_in_executor, _executor,
        _run_pptx_to_word_sync, job_id, in_path, tmp_dir, stem,
    )
    asyncio.create_task(_schedule_job_cleanup(job_id))
    return JSONResponse(status_code=202, content={"job_id": job_id, "status": JobStatus.PENDING})


@router.post("/api/convert/word-to-excel")
async def convert_word_to_excel(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
):
    """Submit an async Word (.docx) to Excel (.xlsx) conversion job."""
    if not file.filename.lower().endswith(('.docx', '.doc')):
        raise HTTPException(status_code=400, detail="Uploaded file must be a Word document (.docx).")

    stem = Path(file.filename).stem
    tmp_dir = tempfile.mkdtemp()
    in_path = os.path.join(tmp_dir, Path(file.filename).name)

    CHUNK = 1024 * 1024
    try:
        with open(in_path, "wb") as out_f:
            while True:
                chunk = await file.read(CHUNK)
                if not chunk:
                    break
                out_f.write(chunk)
    except Exception as e:
        shutil.rmtree(tmp_dir, ignore_errors=True)
        raise HTTPException(status_code=500, detail=f"Upload failed: {e}")

    job_id = str(uuid.uuid4())
    _job_store[job_id] = {
        "status": JobStatus.PENDING, "result_path": None, "tmp_dir": tmp_dir,
        "filename": f"{stem}.xlsx",
        "content_type": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "error": None,
    }
    background_tasks.add_task(
        asyncio.get_event_loop().run_in_executor, _executor,
        _run_word_to_excel_sync, job_id, in_path, tmp_dir, stem,
    )
    asyncio.create_task(_schedule_job_cleanup(job_id))
    return JSONResponse(status_code=202, content={"job_id": job_id, "status": JobStatus.PENDING})


@router.post("/api/convert/excel-to-word")
async def convert_excel_to_word(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
):
    """Submit an async Excel (.xlsx) to Word (.docx) conversion job."""
    if not file.filename.lower().endswith(('.xlsx', '.xls')):
        raise HTTPException(status_code=400, detail="Uploaded file must be an Excel spreadsheet (.xlsx).")

    stem = Path(file.filename).stem
    tmp_dir = tempfile.mkdtemp()
    in_path = os.path.join(tmp_dir, Path(file.filename).name)

    CHUNK = 1024 * 1024
    try:
        with open(in_path, "wb") as out_f:
            while True:
                chunk = await file.read(CHUNK)
                if not chunk:
                    break
                out_f.write(chunk)
    except Exception as e:
        shutil.rmtree(tmp_dir, ignore_errors=True)
        raise HTTPException(status_code=500, detail=f"Upload failed: {e}")

    job_id = str(uuid.uuid4())
    _job_store[job_id] = {
        "status": JobStatus.PENDING, "result_path": None, "tmp_dir": tmp_dir,
        "filename": f"{stem}.docx",
        "content_type": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "error": None,
    }
    background_tasks.add_task(
        asyncio.get_event_loop().run_in_executor, _executor,
        _run_excel_to_word_sync, job_id, in_path, tmp_dir, stem,
    )
    asyncio.create_task(_schedule_job_cleanup(job_id))
    return JSONResponse(status_code=202, content={"job_id": job_id, "status": JobStatus.PENDING})


def _hex_to_rgb(hex_str: str) -> Tuple[float, float, float]:
    """Convert hex color string like '#000000' or 'FF0000' to RGB tuple of floats (0.0 to 1.0)."""
    if not hex_str:
        return (0.0, 0.0, 0.0)
    hex_str = str(hex_str).lstrip('#')
    if len(hex_str) == 6:
        try:
            r = int(hex_str[0:2], 16) / 255.0
            g = int(hex_str[2:4], 16) / 255.0
            b = int(hex_str[4:6], 16) / 255.0
            return (r, g, b)
        except ValueError:
            pass
    return (0.0, 0.0, 0.0)



